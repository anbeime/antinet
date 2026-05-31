import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import pptx
print(f"python-pptx version: {pptx.__version__}")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE as MST

prs = Presentation(r'C:\D\zhiyi\generated\Test_20260531_122758.pptx')

for shape in prs.slides[0].shapes:
    if shape.shape_type == MST.GROUP:
        print(f"GROUP: {shape.name}")
    elif shape.shape_type == 17:
        print(f"TEXT_BOX: {shape.name}")
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                t = r.text[:30].encode('ascii', 'replace').decode('ascii')
                print(f"  text='{t}'")
                print(f"  has _element: {hasattr(r, '_element')}")
                print(f"  font.size: {r.font.size}")
                try:
                    clr = r.font.color.rgb
                    print(f"  font.color.rgb: #{clr}")
                except Exception as e:
                    print(f"  font.color.rgb error: {type(e).__name__}: {e}")
                break
