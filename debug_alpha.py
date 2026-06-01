"""Debug alpha extraction from python-pptx shapes."""
import sys
sys.path.insert(0, 'C:\\D\\zhiyi\\backend')
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

path = Path('C:/D/zhiyi/generated/Test_20260531_122758.pptx')
prs = Presentation(str(path))

slide2 = prs.slides[1]
for shape in slide2.shapes:
    name = shape.name
    stype = shape.shape_type
    if stype == 6:  # GROUP
        print(f"  GROUP {name}")
        for child in shape.shapes:
            cname = child.name
            # Check for alpha
            alpha_el = child._element.find('.//' + qn('a:alpha'))
            if alpha_el is not None:
                val = alpha_el.get('val')
                print(f"    {cname}: alpha found, val={val}")
            # Also check fill
            fill_el = child._element.find('.//' + qn('a:solidFill'))
            if fill_el is not None:
                try:
                    from pptx.enum.dml import MSO_FILL_TYPE
                    print(f"    {cname}: has solidFill, fill.type={child.fill.type}")
                    if child.fill.type == MSO_FILL_TYPE.SOLID:
                        print(f"      color={child.fill.fore_color.rgb}")
                except Exception as e:
                    print(f"    {cname}: fill error: {e}")
