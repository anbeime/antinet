"""
报表自动化服务
完整报表自动化系统:
数据源 → pandas处理 → Excel生成 → PDF报告 → PPT演示 → 网页展示
支持ARM架构
"""

import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Any, Optional
import json

logger = logging.getLogger(__name__)

OUTPUT_DIR = Path("./data/exports")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    logger.warning("pandas 不可用")

try:
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.utils import get_column_letter
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("openpyxl 不可用")

try:
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logger.warning("ReportLab 不可用")

try:
    import sys
    logger.info(f"Python path: {sys.path[:3]}")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PYTHON_PPTX_AVAILABLE = True
    logger.info("python-pptx loaded successfully")
except ImportError as e:
    PYTHON_PPTX_AVAILABLE = False
    import sys
    logger.warning(f"python-pptx 不可用: {e}")
    logger.warning(f"Python path: {sys.path[:3]}")
except Exception as e:
    PYTHON_PPTX_AVAILABLE = False
    logger.warning(f"python-pptx 加载异常: {e}")


class ReportAutomationService:
    """报表自动化服务"""
    
    COLOR_SCHEME = {
        'blue': {'start': '5B9BD5', 'end': '5B9BD5'},
        'green': {'start': '70AD47', 'end': '70AD47'},
        'yellow': {'start': 'FFC000', 'end': 'FFC000'},
        'red': {'start': 'ED7D31', 'end': 'ED7D31'},
        'header': {'start': '4472C4', 'end': '4472C4'},
    }
    
    def __init__(self):
        self.output_dir = OUTPUT_DIR
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
    def process_data(self, data: List[Dict], config: Optional[Dict] = None) -> Dict:
        """使用pandas处理数据"""
        if not PANDAS_AVAILABLE:
            return {'error': 'pandas 不可用', 'data': data}
            
        config = config or {}
        df = pd.DataFrame(data)
        
        stats = {
            'total_rows': len(df),
            'total_columns': len(df.columns),
            'missing_values': int(df.isnull().sum().sum()),
            'duplicates': int(df.duplicated().sum()),
            'columns': []
        }
        
        for col in df.columns:
            col_info = {
                'name': col,
                'type': str(df[col].dtype),
                'missing': int(df[col].isnull().sum()),
                'unique': int(df[col].nunique())
            }
            
            if pd.api.types.is_numeric_dtype(df[col]):
                col_info['stats'] = {
                    'min': float(df[col].min()) if not df[col].isnull().all() else None,
                    'max': float(df[col].max()) if not df[col].isnull().all() else None,
                    'mean': float(df[col].mean()) if not df[col].isnull().all() else None,
                    'median': float(df[col].median()) if not df[col].isnull().all() else None,
                }
            stats['columns'].append(col_info)
            
        return {
            'data': data,
            'stats': stats,
            'summary': self._generate_summary(stats)
        }
        
    def _generate_summary(self, stats: Dict) -> Dict:
        """生成数据摘要"""
        return {
            'fact': f"共处理 {stats['total_rows']} 行数据，包含 {stats['total_columns']} 个字段",
            'interpret': f"数据缺失值 {stats['missing_values']} 个，重复行 {stats['duplicates']} 行",
            'risk': '数据质量良好' if stats['missing_values'] == 0 else '存在数据缺失，建议补充',
            'action': '建议：进行深入数据分析' if stats['missing_values'] == 0 else '建议：清理缺失值'
        }
        
    def generate_excel(
        self,
        data: List[Dict],
        config: Optional[Dict] = None,
        include_charts: bool = True
    ) -> str:
        """生成Excel报表"""
        if not OPENPYXL_AVAILABLE:
            raise RuntimeError("openpyxl 不可用")
            
        config = config or {}
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = config.get('filename') or f'report_{timestamp}.xlsx'
        output_path = self.output_dir / filename
        
        wb = Workbook()
        
        ws_summary = wb.active
        ws_summary.title = "数据摘要"
        self._create_summary_sheet(ws_summary, data, config)
        
        if 'stats' in config:
            ws_stats = wb.create_sheet("统计分析")
            self._create_stats_sheet(ws_stats, config['stats'])
            
        ws_data = wb.create_sheet("原始数据")
        self._create_data_sheet(ws_data, data)
        
        if include_charts and PANDAS_AVAILABLE:
            ws_charts = wb.create_sheet("数据图表")
            self._create_charts_sheet(ws_charts, data)
            
        wb.save(output_path)
        return str(output_path)
        
    def _create_summary_sheet(self, ws, data: List[Dict], config: Dict):
        """创建摘要工作表"""
        ws['A1'] = "报表摘要"
        ws['A1'].font = Font(size=16, bold=True)
        ws.merge_cells('A1:D1')
        
        headers = ['类型', '标题', '内容', '创建时间']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=3, column=col)
            cell.value = header
            cell.font = Font(bold=True, color='FFFFFF')
            cell.fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
            
        cards_by_type = {
            '🔵 事实卡片': config.get('fact', []),
            '🟢 解释卡片': config.get('interpret', []),
            '🟡 风险卡片': config.get('risk', []),
            '🔴 行动建议': config.get('action', [])
        }
        
        row = 4
        for card_type, cards in cards_by_type.items():
            if not cards:
                continue
            for card in cards:
                ws.cell(row=row, column=1).value = card_type
                ws.cell(row=row, column=2).value = card.get('title', '')
                ws.cell(row=row, column=3).value = card.get('content', '')
                ws.cell(row=row, column=4).value = card.get('created_at', datetime.now().isoformat())
                
                color = self.COLOR_SCHEME.get(card_type.split()[0], {}).get('start', 'FFFFFF')
                for col in range(1, 5):
                    ws.cell(row=row, column=col).fill = PatternFill(
                        start_color=color, end_color=color, fill_type='solid'
                    )
                row += 1
                
        for col in range(1, 5):
            ws.column_dimensions[get_column_letter(col)].width = 25
            
    def _create_stats_sheet(self, ws, stats: Dict):
        """创建统计分析工作表"""
        ws['A1'] = "统计分析"
        ws['A1'].font = Font(size=14, bold=True)
        
        ws['A3'] = "指标"
        ws['B3'] = "值"
        for col in ['A', 'B']:
            ws[f'{col}3'].font = Font(bold=True)
            
        metrics = [
            ('总行数', stats.get('total_rows', 0)),
            ('总列数', stats.get('total_columns', 0)),
            ('缺失值', stats.get('missing_values', 0)),
            ('重复行', stats.get('duplicates', 0)),
        ]
        
        row = 4
        for metric, value in metrics:
            ws.cell(row=row, column=1).value = metric
            ws.cell(row=row, column=2).value = value
            row += 1
            
    def _create_data_sheet(self, ws, data: List[Dict]):
        """创建原始数据工作表"""
        if not data:
            return
            
        headers = list(data[0].keys())
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col)
            cell.value = header
            cell.font = Font(bold=True, color='FFFFFF')
            cell.fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
            
        for row_idx, row_data in enumerate(data, 2):
            for col_idx, header in enumerate(headers, 1):
                ws.cell(row=row_idx, column=col_idx).value = row_data.get(header)
                
        for col in range(1, len(headers) + 1):
            ws.column_dimensions[get_column_letter(col)].width = 15
            
    def _create_charts_sheet(self, ws, data: List[Dict]):
        """创建图表工作表"""
        if not data:
            return
            
        numeric_cols = []
        for key in data[0].keys():
            if isinstance(data[0].get(key), (int, float)):
                numeric_cols.append(key)
                
        if not numeric_cols:
            return
            
        ws['A1'] = "数据图表"
        ws['A1'].font = Font(size=14, bold=True)
        
        if len(numeric_cols) >= 2:
            chart = BarChart()
            chart.title = "数值对比"
            chart.style = 10
            
            ws['A3'] = "字段"
            ws['B3'] = "值"
            for col in ['A', 'B']:
                ws[f'{col}3'].font = Font(bold=True)
                
            row = 4
            for col in numeric_cols[:5]:
                values = [d.get(col, 0) for d in data[:10]]
                ws.cell(row=row, column=1).value = col
                ws.cell(row=row, column=2).value = sum(values) / len(values) if values else 0
                row += 1
                
            data_ref = Reference(ws, min_col=2, min_row=3, max_row=row-1)
            categories = Reference(ws, min_col=1, min_row=4, max_row=row-1)
            chart.add_data(data_ref, titles_from_data=True)
            chart.set_categories(categories)
            ws.add_chart(chart)
            
    def _get_chinese_font_path(self) -> Optional[str]:
        """查找中文字体路径"""
        import os
        # 多个可能的字体路径
        base_dirs = [
            Path(__file__).parent.parent.parent,  # 项目根目录
            Path(__file__).parent.parent,  # backend 目录
        ]
        font_names = [
            "NotoSansSC-Regular.ttf",
            "NotoSansSC-Bold.ttf",
            "SimHei.ttf",
            "msyh.ttc",
        ]
        for base in base_dirs:
            for font_name in font_names:
                # 尝试 public/fonts, fonts, static/fonts 等目录
                for subdir in ["public/fonts", "fonts", "static/fonts", "assets/fonts"]:
                    path = base / subdir / font_name
                    if path.exists():
                        return str(path)
                # 直接在 base 下查找
                path = base / font_name
                if path.exists():
                    return str(path)
        return None

    def _register_chinese_font(self):
        """注册中文字体（只执行一次）"""
        if not hasattr(self, '_font_registered'):
            font_path = self._get_chinese_font_path()
            if font_path:
                try:
                    pdfmetrics.registerFont(TTFont('ChineseFont', font_path, 'Identity-H'))
                    self._font_registered = True
                    self._font_name = 'ChineseFont'
                    logger.info(f"[PDF] 中文字体注册成功: {font_path}")
                except Exception as e:
                    logger.warning(f"[PDF] 中文字体注册失败: {e}")
                    self._font_registered = True
                    self._font_name = 'Helvetica'
            else:
                logger.warning("[PDF] 未找到中文字体，PDF中文可能显示乱码")
                self._font_registered = True
                self._font_name = 'Helvetica'

    def generate_pdf(
        self,
        data: Dict,
        config: Optional[Dict] = None
    ) -> str:
        """生成PDF报表"""
        if not REPORTLAB_AVAILABLE:
            raise RuntimeError("ReportLab 不可用")

        # 注册中文字体
        self._register_chinese_font()
        font_name = getattr(self, '_font_name', 'Helvetica')

        config = config or {}
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = config.get('filename') or f'report_{timestamp}.pdf'
        output_path = self.output_dir / filename

        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=A4,
            rightMargin=72, leftMargin=72,
            topMargin=72, bottomMargin=18
        )

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontName=font_name,
            fontSize=18,
            spaceAfter=30
        )
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontName=font_name,
            fontSize=14,
            spaceAfter=12
        )
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontName=font_name,
            fontSize=10,
            spaceAfter=6
        )
        
        story = []
        
        title = config.get('title', '报表分析报告')
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.2*inch))
        
        story.append(Paragraph(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", normal_style))
        story.append(Spacer(1, 0.3*inch))

        cards_by_type = {
            '事实卡片': data.get('fact', []),
            '解释卡片': data.get('interpret', []),
            '风险卡片': data.get('risk', []),
            '行动建议': data.get('action', [])
        }

        for card_type, cards in cards_by_type.items():
            if not cards:
                continue
            story.append(Paragraph(card_type, heading_style))

            for card in cards:
                story.append(Paragraph(f"<b>{card.get('title', '')}</b>", normal_style))
                story.append(Paragraph(card.get('content', ''), normal_style))
                story.append(Spacer(1, 0.1*inch))

        if 'stats' in data:
            story.append(PageBreak())
            story.append(Paragraph("统计分析", heading_style))

            stats = data['stats']
            stats_data = [
                ['指标', '值'],
                ['总行数', str(stats.get('total_rows', 0))],
                ['总列数', str(stats.get('total_columns', 0))],
                ['缺失值', str(stats.get('missing_values', 0))],
                ['重复行', str(stats.get('duplicates', 0))],
            ]

            table = Table(stats_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), font_name + '-Bold' if font_name == 'ChineseFont' else 'Helvetica-Bold'),
                ('FONTNAME', (0, 1), (-1, -1), font_name),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            story.append(table)
            
        doc.build(story)
        return str(output_path)
        
    def generate_ppt(
        self,
        data: Dict,
        config: Optional[Dict] = None
    ) -> str:
        """生成PPT演示文稿"""
        if not PYTHON_PPTX_AVAILABLE:
            raise RuntimeError("python-pptx 不可用")
            
        config = config or {}
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = config.get('filename') or f'report_{timestamp}.pptx'
        output_path = self.output_dir / filename
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = config.get('title', '报表分析报告')
        subtitle.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        bullet_slide_layout = prs.slide_layouts[1]
        
        cards_by_type = {
            '🔵 事实卡片': data.get('fact', []),
            '🟢 解释卡片': data.get('interpret', []),
            '🟡 风险卡片': data.get('risk', []),
            '🔴 行动建议': data.get('action', [])
        }
        
        for card_type, cards in cards_by_type.items():
            if not cards:
                continue
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = card_type
            
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            
            for card in cards:
                p = tf.add_paragraph()
                p.text = card.get('title', '')
                p.level = 0
                
                p = tf.add_paragraph()
                p.text = card.get('content', '')
                p.level = 1
                
        if 'stats' in data:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = "统计分析"
            
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            
            stats = data['stats']
            for key, value in [
                ('总行数', stats.get('total_rows', 0)),
                ('总列数', stats.get('total_columns', 0)),
                ('缺失值', stats.get('missing_values', 0)),
                ('重复行', stats.get('duplicates', 0)),
            ]:
                p = tf.add_paragraph()
                p.text = f"{key}: {value}"
                
        prs.save(str(output_path))
        return str(output_path)
        
    def generate_full_report(
        self,
        data: List[Dict],
        config: Optional[Dict] = None
    ) -> Dict:
        """生成完整报表（Excel + PDF + PPT）"""
        config = config or {}
        
        processed = self.process_data(data, config)
        cards_data = self._generate_cards_from_data(processed)
        
        excel_file = self.generate_excel(
            data,
            config={**config, **cards_data}
        )
        
        pdf_file = self.generate_pdf(cards_data, config)
        
        ppt_file = self.generate_ppt(cards_data, config)
        
        return {
            'excel': excel_file,
            'pdf': pdf_file,
            'ppt': ppt_file,
            'timestamp': datetime.now().isoformat()
        }
        
    def _generate_cards_from_data(self, processed: Dict) -> Dict:
        """从处理后的数据生成4色卡片"""
        stats = processed.get('stats', {})
        
        return {
            'fact': [{
                'title': '数据概览',
                'content': processed['summary']['fact'],
                'created_at': datetime.now().isoformat()
            }],
            'interpret': [{
                'title': '数据质量分析',
                'content': processed['summary']['interpret'],
                'created_at': datetime.now().isoformat()
            }],
            'risk': [{
                'title': '风险提示',
                'content': processed['summary']['risk'],
                'created_at': datetime.now().isoformat()
            }],
            'action': [{
                'title': '行动建议',
                'content': processed['summary']['action'],
                'created_at': datetime.now().isoformat()
            }],
            'stats': stats
        }


service = ReportAutomationService()