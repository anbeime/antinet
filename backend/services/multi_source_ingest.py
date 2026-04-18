"""
多源数据接入模块 - 通政司扩展
支持多格式文档解析、上下文捕获、跨平台数据同步
"""
import logging
import json
import re
import os
from typing import Dict, List, Optional, Any
from datetime import datetime
from pathlib import Path
import hashlib
import mimetypes
import base64

logger = logging.getLogger(__name__)


class MultiSourceIngest:
    """多源数据接入器 - 通政司扩展"""
    
    def __init__(self, storage_path: str = None):
        self.storage_path = storage_path or "./data/ingested"
        Path(self.storage_path).mkdir(parents=True, exist_ok=True)
        self.supported_formats = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.doc': self._parse_docx,
            '.xlsx': self._parse_excel,
            '.xls': self._parse_excel,
            '.pptx': self._parse_pptx,
            '.ppt': self._parse_pptx,
            '.txt': self._parse_text,
            '.md': self._parse_text,
            '.json': self._parse_json,
            '.csv': self._parse_csv,
            '.png': self._parse_image,
            '.jpg': self._parse_image,
            '.jpeg': self._parse_image,
            '.gif': self._parse_image,
        }
    
    def ingest(self, source: str, content: Any = None, file_path: str = None,
               metadata: Dict = None) -> Dict:
        """
        接入数据
        
        参数:
            source: 数据来源类型 (email/im/local/cloud/bookmark)
            content: 内容（文本/二进制）
            file_path: 文件路径（可选）
            metadata: 元数据
        
        返回:
            接入结果
        """
        try:
            logger.info(f"[Ingest] 从 {source} 接入数据")
            
            # 解析内容
            parsed_data = self._parse_content(source, content, file_path)
            
            # 提取元数据
            extracted_metadata = self._extract_metadata(source, content, file_path, metadata)
            
            # 构建标准化数据
            standardized = {
                "content": parsed_data["content"],
                "format": parsed_data["format"],
                "metadata": extracted_metadata,
                "timestamp": datetime.now().isoformat(),
                "content_hash": self._compute_hash(parsed_data["content"])
            }
            
            # 存储
            doc_id = self._save_document(standardized)
            
            return {
                "status": "success",
                "doc_id": doc_id,
                "format": parsed_data["format"],
                "metadata": extracted_metadata
            }
            
        except Exception as e:
            logger.error(f"[Ingest] 接入失败: {e}")
            return {"status": "error", "error": str(e)}
    
    def _parse_content(self, source: str, content: Any, file_path: str) -> Dict:
        """解析内容"""
        # 文件路径解析
        if file_path:
            ext = Path(file_path).suffix.lower()
            if ext in self.supported_formats:
                return self.supported_formats[ext](file_path)
        
        # 直接内容解析
        if isinstance(content, str):
            return {"content": content, "format": "text"}
        elif isinstance(content, bytes):
            return {"content": base64.b64encode(content).decode(), "format": "binary"}
        elif isinstance(content, dict):
            return {"content": json.dumps(content, ensure_ascii=False), "format": "json"}
        else:
            return {"content": str(content), "format": "text"}
    
    def _parse_pdf(self, file_path: str) -> Dict:
        """解析PDF"""
        try:
            import PyPDF2
            text = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text.append(page.extract_text())
            return {"content": "\n".join(text), "format": "pdf"}
        except Exception as e:
            logger.warning(f"PDF解析失败: {e}")
            return {"content": f"[PDF文件: {file_path}]", "format": "pdf"}
    
    def _parse_docx(self, file_path: str) -> Dict:
        """解析Word"""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return {"content": "\n".join(paragraphs), "format": "docx"}
        except Exception as e:
            logger.warning(f"Word解析失败: {e}")
            return {"content": f"[Word文件: {file_path}]", "format": "docx"}
    
    def _parse_excel(self, file_path: str) -> Dict:
        """解析Excel"""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            sheets_data = {}
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    rows.append([str(cell) if cell is not None else "" for cell in row])
                sheets_data[sheet_name] = rows
            
            return {"content": json.dumps(sheets_data, ensure_ascii=False), "format": "xlsx"}
        except Exception as e:
            logger.warning(f"Excel解析失败: {e}")
            return {"content": f"[Excel文件: {file_path}]", "format": "xlsx"}
    
    def _parse_pptx(self, file_path: str) -> Dict:
        """解析PPT"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            
            return {"content": "\n".join(slides_text), "format": "pptx"}
        except Exception as e:
            logger.warning(f"PPT解析失败: {e}")
            return {"content": f"[PPT文件: {file_path}]", "format": "pptx"}
    
    def _parse_text(self, file_path: str) -> Dict:
        """解析文本"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return {"content": content, "format": "text"}
    
    def _parse_json(self, file_path: str) -> Dict:
        """解析JSON"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return {"content": json.dumps(data, ensure_ascii=False), "format": "json"}
    
    def _parse_csv(self, file_path: str) -> Dict:
        """解析CSV"""
        import csv
        rows = []
        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        return {"content": json.dumps(rows, ensure_ascii=False), "format": "csv"}
    
    def _parse_image(self, file_path: str) -> Dict:
        """解析图片 - 记录路径和基本信息"""
        file_size = os.path.getsize(file_path)
        return {
            "content": f"[图片文件: {file_path}, 大小: {file_size} bytes]",
            "format": "image",
            "image_path": file_path
        }
    
    def _extract_metadata(self, source: str, content: Any, file_path: str, 
                          provided_metadata: Dict = None) -> Dict:
        """提取元数据"""
        metadata = provided_metadata or {}
        
        # 来源信息
        metadata["source"] = source
        metadata["ingested_at"] = datetime.now().isoformat()
        
        # 文件信息
        if file_path:
            path = Path(file_path)
            metadata["file_name"] = path.name
            metadata["file_size"] = path.stat().st_size
            metadata["file_extension"] = path.suffix
            
            # 文件时间信息
            stat = path.stat()
            metadata["created_time"] = datetime.fromtimestamp(stat.st_ctime).isoformat()
            metadata["modified_time"] = datetime.fromtimestamp(stat.st_mtime).isoformat()
        
        # 文本内容特征
        if isinstance(content, str):
            metadata["char_count"] = len(content)
            metadata["line_count"] = len(content.split('\n'))
            
            # 提取关键词
            words = re.findall(r'\b\w{2,}\b', content)
            word_freq = {}
            for word in words:
                word_freq[word] = word_freq.get(word, 0) + 1
            top_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:10]
            metadata["keywords"] = [w[0] for w in top_words]
        
        return metadata
    
    def _save_document(self, document: Dict) -> str:
        """保存文档"""
        doc_id = f"doc_{document['content_hash'][:16]}"
        
        doc_path = Path(self.storage_path) / f"{doc_id}.json"
        with open(doc_path, 'w', encoding='utf-8') as f:
            json.dump(document, f, ensure_ascii=False, indent=2)
        
        logger.info(f"[Ingest] 文档已保存: {doc_id}")
        return doc_id
    
    def _compute_hash(self, content: str) -> str:
        """计算内容哈希"""
        return hashlib.sha256(content.encode()).hexdigest()


class EmailIngest(MultiSourceIngest):
    """邮件数据接入"""
    
    def ingest_email(self, email_data: Dict) -> Dict:
        """接入邮件数据"""
        content_parts = []
        
        # 邮件头
        subject = email_data.get("subject", "无主题")
        sender = email_data.get("from", "未知发件人")
        date = email_data.get("date", "")
        
        content_parts.append(f"主题: {subject}")
        content_parts.append(f"发件人: {sender}")
        content_parts.append(f"日期: {date}")
        
        # 邮件正文
        body = email_data.get("body", "")
        content_parts.append(f"\n正文:\n{body}")
        
        # 附件信息
        attachments = email_data.get("attachments", [])
        if attachments:
            content_parts.append(f"\n附件: {', '.join([a.get('name', 'unknown') for a in attachments])}")
        
        return self.ingest(
            source="email",
            content="\n".join(content_parts),
            metadata={
                "email_subject": subject,
                "email_from": sender,
                "email_date": date,
                "has_attachments": len(attachments) > 0
            }
        )


class IMIngest(MultiSourceIngest):
    """即时通讯数据接入"""
    
    def ingest_im_message(self, message_data: Dict) -> Dict:
        """接入IM消息"""
        sender = message_data.get("sender", "未知")
        content = message_data.get("content", "")
        timestamp = message_data.get("timestamp", "")
        
        return self.ingest(
            source="im",
            content=content,
            metadata={
                "im_sender": sender,
                "im_timestamp": timestamp,
                "im_platform": message_data.get("platform", "unknown")
            }
        )


class LocalDocIngest(MultiSourceIngest):
    """本地文档接入 - 扫描目录"""
    
    def scan_directory(self, directory: str, extensions: List[str] = None) -> List[Dict]:
        """扫描目录"""
        if extensions is None:
            extensions = list(self.supported_formats.keys())
        
        results = []
        dir_path = Path(directory)
        
        for ext in extensions:
            for file_path in dir_path.rglob(f"*{ext}"):
                try:
                    result = self.ingest(
                        source="local",
                        file_path=str(file_path),
                        metadata={"scan_directory": directory}
                    )
                    results.append(result)
                except Exception as e:
                    logger.warning(f"处理文件失败 {file_path}: {e}")
        
        return results


class BookmarkIngest(MultiSourceIngest):
    """浏览器书签接入"""
    
    def ingest_bookmarks(self, bookmarks_data: List[Dict]) -> List[Dict]:
        """接入书签数据"""
        results = []
        
        for bookmark in bookmarks_data:
            content = f"标题: {bookmark.get('title', '')}\nURL: {bookmark.get('url', '')}"
            
            if bookmark.get("children"):
                for child in bookmark["children"]:
                    child_content = f"标题: {child.get('title', '')}\nURL: {child.get('url', '')}"
                    result = self.ingest(
                        source="bookmark",
                        content=child_content,
                        metadata={
                            "folder": bookmark.get("title", ""),
                            "bookmark_title": child.get("title", "")
                        }
                    )
                    results.append(result)
            else:
                result = self.ingest(
                    source="bookmark",
                    content=content,
                    metadata={"bookmark_title": bookmark.get("title", "")}
                )
                results.append(result)
        
        return results


class CloudNoteIngest(MultiSourceIngest):
    """云端笔记接入 (OneNote, 印象笔记等)"""
    
    async def ingest_onenote(self, note_data: Dict) -> Dict:
        """接入OneNote笔记"""
        content_parts = []
        
        title = note_data.get("title", "")
        content_parts.append(f"标题: {title}")
        
        sections = note_data.get("sections", [])
        for section in sections:
            content_parts.append(f"\n## {section.get('name', '')}")
            pages = section.get("pages", [])
            for page in pages:
                content_parts.append(f"\n### {page.get('title', '')}")
                content_parts.append(page.get("content", ""))
        
        return self.ingest(
            source="onenote",
            content="\n".join(content_parts),
            metadata={
                "notebook": note_data.get("notebook", ""),
                "page_count": sum(len(s.get("pages", [])) for s in sections)
            }
        )
    
    async def ingest_evernote(self, note_data: Dict) -> Dict:
        """接入印象笔记"""
        content_parts = []
        
        title = note_data.get("title", "")
        content_parts.append(f"标题: {title}")
        
        content_parts.append(f"\n内容:\n{note_data.get('content', '')}")
        
        tags = note_data.get("tags", [])
        if tags:
            content_parts.append(f"\n标签: {', '.join(tags)}")
        
        return self.ingest(
            source="evernote",
            content="\n".join(content_parts),
            metadata={
                "evernote_guid": note_data.get("guid", ""),
                "tags": tags
            }
        )


class ContextCapture:
    """上下文捕获器 - 自动记录上下文信息"""
    
    def __init__(self):
        self.context_stack = []
    
    def capture_context(self, event_type: str, event_data: Dict) -> Dict:
        """捕获上下文"""
        context = {
            "event_type": event_type,
            "timestamp": datetime.now().isoformat(),
            "event_data": event_data,
            "context_stack": self.context_stack.copy()
        }
        
        self.context_stack.append({
            "type": event_type,
            "timestamp": context["timestamp"]
        })
        
        # 保持最近20条上下文
        if len(self.context_stack) > 20:
            self.context_stack = self.context_stack[-20:]
        
        return context
    
    def get_current_context(self) -> Dict:
        """获取当前上下文"""
        if not self.context_stack:
            return {}
        
        return {
            "recent_events": self.context_stack[-5:],
            "event_count": len(self.context_stack)
        }
    
    def clear_context(self):
        """清空上下文"""
        self.context_stack = []