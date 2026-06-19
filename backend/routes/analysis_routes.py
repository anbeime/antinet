"""
多类型文件分析路由
支持 Excel, JSON, Markdown, PPT 等文件类型的分析和提取
"""
from fastapi import APIRouter, UploadFile, File, HTTPException
from typing import Dict, Any
import io
import json
import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/analysis", tags=["analysis"])


# ============ Excel 分析 ============
async def analyze_excel(contents: bytes, filename: str) -> Dict[str, Any]:
    """分析 Excel 文件"""
    import pandas as pd
    try:
        df = pd.read_excel(io.BytesIO(contents))
    except Exception as e:
        logger.error(f"[AnalysisRoutes] Excel读取失败: {e}")
        raise HTTPException(status_code=400, detail=f"Excel文件读取失败: {str(e)}")

    logger.info(f"[AnalysisRoutes] Excel数据形状: {df.shape}")

    # 分析列信息
    columns = []
    for col in df.columns:
        dtype = str(df[col].dtype)
        col_type = 'string'
        if 'int' in dtype or 'float' in dtype:
            col_type = 'number'
        elif 'datetime' in dtype:
            col_type = 'date'
        elif 'bool' in dtype:
            col_type = 'boolean'

        sample = None
        if len(df) > 0:
            sample_val = df[col].iloc[0]
            if pd.notna(sample_val):
                sample = float(sample_val) if col_type == 'number' else str(sample_val)

        columns.append({
            'key': str(col),
            'name': str(col),
            'type': col_type,
            'sample': sample
        })

    stats = {
        'totalRows': len(df),
        'totalColumns': len(df.columns),
        'numericColumns': len(df.select_dtypes(include=['number']).columns),
        'textColumns': len(df.select_dtypes(include=['object']).columns),
        'dateColumns': len(df.select_dtypes(include=['datetime']).columns),
        'missingValues': int(df.isnull().sum().sum()),
        'duplicates': int(df.duplicated().sum())
    }

    max_rows = min(1000, len(df))
    data = df.head(max_rows).fillna('').to_dict('records')

    for row in data:
        for key, value in row.items():
            if pd.isna(value):
                row[key] = None
            elif isinstance(value, (pd.Timestamp, pd.DatetimeTZDtype)):
                row[key] = str(value)

    return {
        'success': True,
        'fileType': 'excel',
        'filename': filename,
        'data': data,
        'columns': columns,
        'stats': stats,
        'message': f'成功分析 Excel: {stats["totalRows"]} 行 x {stats["totalColumns"]} 列'
    }


# ============ JSON 分析 ============
def analyze_json(contents: bytes, filename: str) -> Dict[str, Any]:
    """分析 JSON 文件"""
    try:
        data = json.loads(contents)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=400, detail=f"JSON格式错误: {str(e)}")

    def get_structure(obj, depth=0, max_depth=10):
        """递归分析 JSON 结构"""
        if depth > max_depth:
            return {'type': 'max_depth_exceeded'}
        if isinstance(obj, dict):
            return {
                'type': 'object',
                'keys': list(obj.keys())[:20],  # 最多20个key
                'key_count': len(obj.keys()),
                'children': {k: get_structure(v, depth+1) for k, v in list(obj.items())[:10]}
            }
        elif isinstance(obj, list):
            sample = obj[:3] if len(obj) > 0 else []
            return {
                'type': 'array',
                'length': len(obj),
                'item_type': get_structure(obj[0], depth+1) if obj else {'type': 'empty'},
                'sample': [get_structure(s, depth+1) for s in sample]
            }
        else:
            return {'type': type(obj).__name__, 'sample': str(obj)[:100] if obj is not None else None}

    def count_stats(obj, _stats=None):
        """统计 JSON 数据"""
        if _stats is None:
            _stats = {'total_keys': 0, 'max_depth': 0, 'lists': 0, 'objects': 0, 'primitives': 0}
        if isinstance(obj, dict):
            _stats['objects'] += 1
            _stats['total_keys'] += len(obj)
            for v in obj.values():
                count_stats(v, _stats)
        elif isinstance(obj, list):
            _stats['lists'] += 1
            for item in obj:
                count_stats(item, _stats)
        else:
            _stats['primitives'] += 1
        return _stats

    structure = get_structure(data)
    stats = count_stats(data)
    stats['filename'] = filename

    # 提取顶层 keys
    top_keys = list(data.keys()) if isinstance(data, dict) else []
    if isinstance(data, list) and len(data) > 0:
        top_keys = list(data[0].keys()) if isinstance(data[0], dict) else ['[array items]']

    return {
        'success': True,
        'fileType': 'json',
        'filename': filename,
        'topKeys': top_keys,
        'structure': structure,
        'stats': stats,
        'preview': data if isinstance(data, dict) else (data[:5] if isinstance(data, list) else data),
        'message': f'JSON 解析成功: 深度 {stats.get("max_depth", 0)}, {stats["total_keys"]} 个键值'
    }


# ============ Markdown 分析 ============
def analyze_markdown(contents: bytes, filename: str) -> Dict[str, Any]:
    """分析 Markdown 文件"""
    text = contents.decode('utf-8', errors='replace')

    # 统计
    lines = text.split('\n')
    total_lines = len(lines)
    total_chars = len(text)
    total_words = len(re.findall(r'\S+', text))

    # 提取标题结构
    headings = []
    for i, line in enumerate(lines):
        match = re.match(r'^(#{1,6})\s+(.+)', line)
        if match:
            level = len(match.group(1))
            title = match.group(2).strip()
            headings.append({
                'level': level,
                'title': title,
                'line': i + 1
            })

    # 提取代码块
    code_blocks = []
    in_code = False
    code_start = 0
    code_content = []
    code_lang = ''
    for i, line in enumerate(lines):
        if line.strip().startswith('```'):
            if not in_code:
                in_code = True
                code_start = i
                code_lang = line.strip()[3:].strip()
            else:
                in_code = False
                code_blocks.append({
                    'language': code_lang,
                    'lines': code_start + 1,
                    'content': '\n'.join(code_content)[:500]
                })
                code_content = []
                code_lang = ''
        elif in_code:
            code_content.append(line)

    # 提取列表项
    list_items = []
    for i, line in enumerate(lines):
        match = re.match(r'^(\s*)[-*+]\s+(.+)', line)
        if match:
            list_items.append({
                'indent': len(match.group(1)),
                'text': match.group(2).strip(),
                'line': i + 1
            })

    # 提取表格
    tables = []
    table_start = -1
    table_lines = []
    for i, line in enumerate(lines):
        if '|' in line:
            if table_start == -1:
                table_start = i
            table_lines.append(line)
        else:
            if table_start != -1 and len(table_lines) >= 2:
                tables.append({
                    'startLine': table_start + 1,
                    'rows': len(table_lines)
                })
            table_start = -1
            table_lines = []
    if table_start != -1 and len(table_lines) >= 2:
        tables.append({'startLine': table_start + 1, 'rows': len(table_lines)})

    # 提取链接和图片
    links = re.findall(r'\[([^\]]+)\]\(([^\)]+)\)', text)

    stats = {
        'totalLines': total_lines,
        'totalChars': total_chars,
        'totalWords': total_words,
        'headings': len(headings),
        'codeBlocks': len(code_blocks),
        'listItems': len(list_items),
        'tables': len(tables),
        'links': len(links),
        'filename': filename
    }

    return {
        'success': True,
        'fileType': 'markdown',
        'filename': filename,
        'headings': headings[:30],  # 最多30个标题
        'codeBlocks': code_blocks[:10],  # 最多10个代码块
        'listItems': list_items[:50],
        'tables': tables[:10],
        'links': links[:20],
        'stats': stats,
        'preview': text[:1000],
        'message': f'Markdown 解析成功: {len(headings)} 个标题, {total_words} 个词'
    }


# ============ PPT 分析 ============
def analyze_pptx(contents: bytes, filename: str) -> Dict[str, Any]:
    """分析 PPT 文件"""
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx 库未安装")

    try:
        prs = Presentation(io.BytesIO(contents))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPT 文件读取失败: {str(e)}")

    slides = []
    for i, slide in enumerate(prs.slides):
        slide_data = {
            'index': i + 1,
            'title': '',
            'content': [],
            'shapes': len(slide.shapes)
        }
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # 检查是否是标题
                    if hasattr(shape, 'text_frame') and shape.text_frame.style:
                        style_name = getattr(shape.text_frame, 'style', '') or ''
                        if 'title' in style_name.lower() or 'heading' in style_name.lower():
                            if not slide_data['title']:
                                slide_data['title'] = text[:100]
                        else:
                            if text not in slide_data['content']:
                                slide_data['content'].append(text[:200])
                    else:
                        if text not in slide_data['content']:
                            slide_data['content'].append(text[:200])

        slides.append(slide_data)

    stats = {
        'totalSlides': len(slides),
        'filename': filename,
        'slideLayouts': list(set(s.layout.name for s in prs.slides if s.layout))
    }

    return {
        'success': True,
        'fileType': 'pptx',
        'filename': filename,
        'slides': slides,
        'stats': stats,
        'message': f'PPT 解析成功: {len(slides)} 张幻灯片'
    }


# ============ PDF 分析 ============
def analyze_pdf(contents: bytes, filename: str) -> Dict[str, Any]:
    """分析 PDF 文件 - 提取文本和基本信息"""
    try:
        from pypdf import PdfReader
    except ImportError:
        raise HTTPException(status_code=500, detail="pypdf 库未安装")

    try:
        reader = PdfReader(io.BytesIO(contents))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF 文件读取失败: {str(e)}")

    pages = []
    total_chars = 0
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ''
        total_chars += len(text)
        pages.append({
            'index': i + 1,
            'text': text[:2000],
            'chars': len(text)
        })

    stats = {
        'totalPages': len(reader.pages),
        'totalChars': total_chars,
        'filename': filename
    }

    return {
        'success': True,
        'fileType': 'pdf',
        'filename': filename,
        'pages': pages[:20],
        'stats': stats,
        'message': f'PDF 解析成功: {len(reader.pages)} 页'
    }


# ============ 主入口 ============
@router.post("/upload-and-analyze")
async def upload_and_analyze(file: UploadFile = File(...)):
    """
    上传并分析文件

    支持的文件类型:
    - Excel: .xlsx, .xls
    - JSON: .json
    - Markdown: .md
    - PPT: .pptx
    - PDF: .pdf
    """
    logger.info(f"[AnalysisRoutes] 收到文件上传: {file.filename}")

    try:
        contents = await file.read()
        filename = file.filename

        # 根据文件扩展名路由到对应分析器
        if filename.endswith(('.xlsx', '.xls')):
            return await analyze_excel(contents, filename)
        elif filename.endswith('.json'):
            return analyze_json(contents, filename)
        elif filename.endswith('.md'):
            return analyze_markdown(contents, filename)
        elif filename.endswith('.pptx'):
            return analyze_pptx(contents, filename)
        elif filename.endswith('.pdf'):
            return analyze_pdf(contents, filename)
        else:
            raise HTTPException(
                status_code=400,
                detail=f"不支持的文件格式: {Path(filename).suffix}\n支持的格式: .xlsx, .xls, .json, .md, .pptx, .pdf"
            )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[AnalysisRoutes] 分析失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"分析失败: {str(e)}")


@router.get("/health")
async def health_check():
    """健康检查"""
    return {
        "status": "healthy",
        "service": "analysis",
        "supported_formats": [".xlsx", ".xls", ".json", ".md", ".pptx", ".pdf"]
    }