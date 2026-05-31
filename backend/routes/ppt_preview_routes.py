"""
PPT 预览路由（增强版）
- 提取完整 PPT 样式（颜色、字体、形状位置）
- 支持导入时保持原始样式还原
- 支持编辑后重新生成
"""
import logging
import io
import base64
import json
from pathlib import Path
from fastapi import APIRouter, HTTPException, UploadFile, File, Body
from fastapi.responses import JSONResponse
from typing import List, Dict, Any, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx 未安装")

from design_system import (
    DesignPresets, BrandStyle, SlideShape, SlideData, PPTPreviewData,
    rgb_to_hex, hex_to_rgb, LayoutStyle, ThemeColors, ThemeFonts,
)

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api/ppt", tags=["PPT Preview (Enhanced)"])


def _extract_shapes_recursive(shape) -> list[Dict[str, Any]]:
    """Recursively extract shapes from a shape (handles group shapes).

    python-pptx slide.shapes only returns top-level shapes, not children of
    group shapes (<p:grpSp>). This function recursively walks into groups
    and returns a flat list of all leaf shapes with text/table/fill.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    results: list[Dict[str, Any]] = []

    # If this is a GROUP shape, recurse into its children
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        try:
            for child in shape.shapes:
                results.extend(_extract_shapes_recursive(child))
        except Exception:
            pass
        return results

    data = _extract_shape_style(shape)
    if data.get("paragraphs") or data.get("table") or data.get("fill_color"):
        results.append(data)
    return results


def _extract_shape_style(shape) -> Dict[str, Any]:
    """提取单个形状的完整样式（不处理 group 嵌套）"""
    data = {
        "type": str(shape.shape_type),
        "name": shape.name,
        "left": round(shape.left.emu / 914400 * 96, 2) if hasattr(shape, 'left') else 0,
        "top": round(shape.top.emu / 914400 * 96, 2) if hasattr(shape, 'top') else 0,
        "width": round(shape.width.emu / 914400 * 96, 2) if hasattr(shape, 'width') else 0,
        "height": round(shape.height.emu / 914400 * 96, 2) if hasattr(shape, 'height') else 0,
    }

    if shape.has_text_frame:
        paragraphs = []
        full_text = ""
        for para in shape.text_frame.paragraphs:
            runs_data = []
            para_text = ""
            for run in para.runs:
                run_info = {"text": run.text}
                try:
                    if run.font.size:
                        run_info["font_size"] = round(run.font.size.pt, 1)
                    if run.font.bold:
                        run_info["bold"] = True
                    if run.font.italic:
                        run_info["italic"] = True
                    if run.font.color and run.font.color.rgb:
                        run_info["color"] = f"#{run.font.color.rgb}"
                    if run.font.name:
                        run_info["font_name"] = run.font.name
                except:
                    pass
                runs_data.append(run_info)
                para_text += run.text

            para_info = {"text": para_text, "runs": runs_data}
            try:
                align_map = {
                    PP_ALIGN.LEFT: "left",
                    PP_ALIGN.CENTER: "center",
                    PP_ALIGN.RIGHT: "right",
                }
                if para.alignment:
                    para_info["align"] = align_map.get(para.alignment, "left")
            except:
                pass
            paragraphs.append(para_info)
            full_text += para_text

        data["paragraphs"] = paragraphs

        if full_text:
            first_run = None
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    first_run = run
                    break
                if first_run:
                    break
            if first_run:
                try:
                    if first_run.font.size:
                        data["default_font_size"] = round(first_run.font.size.pt, 1)
                        data["font_size"] = round(first_run.font.size.pt, 1)
                    if first_run.font.color and first_run.font.color.rgb:
                        clr = f"#{first_run.font.color.rgb}"
                        data["default_font_color"] = clr
                        data["font_color"] = clr
                    if first_run.font.name:
                        data["default_font_name"] = first_run.font.name
                except:
                    pass

    # Extract fill alpha (opacity) from a:alpha element
    try:
        from pptx.oxml.ns import qn
        alpha_el = shape._element.find('.//' + qn('a:alpha'))
        if alpha_el is not None:
            val = alpha_el.get('val')
            if val:
                alpha_pct = int(val) / 100000.0
                if alpha_pct < 1.0:
                    data["fill_opacity"] = round(alpha_pct, 4)
    except Exception:
        pass

    if shape.has_table:
        table_data = []
        for row in shape.table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            table_data.append(row_data)
        data["table"] = table_data

    try:
        if hasattr(shape, 'fill') and shape.fill.type is not None:
            if shape.fill.type == 1:
                color = shape.fill.fore_color.rgb
                data["fill_color"] = f"#{color}"
    except:
        pass

    return data


def _detect_theme_from_pptx(prs: Presentation) -> str:
    """从 PPTX 内容自动检测最匹配的设计主题"""
    colors_found = {"primary": None, "background": None, "text": None}

    for slide in prs.slides:
        try:
            if slide.background.fill.type is not None and slide.background.fill.type == 1:
                bg = slide.background.fill.fore_color.rgb
                colors_found["background"] = f"#{bg}"
        except:
            pass

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size.pt >= 24:
                            try:
                                if run.font.color and run.font.color.rgb:
                                    colors_found["primary"] = f"#{run.font.color.rgb}"
                            except:
                                pass

    best = "professional"
    if colors_found["background"]:
        bg_lower = colors_found["background"].lower()
        if bg_lower in ("#ffffff", "#f8f9fa", "#f8fafc"):
            best = "minimal"
        elif bg_lower in ("#0f172a", "#1a1a2e", "#1e293b"):
            best = "tech"
    return best


@router.get("/preview/status")
async def get_preview_status():
    return {
        "available": PPTX_AVAILABLE,
        "message": "PPT预览功能已就绪" if PPTX_AVAILABLE else "请安装: pip install python-pptx"
    }


@router.post("/preview/extract")
async def extract_ppt_content(file: UploadFile = File(...)):
    """提取 PPT 完整样式数据，保留原始设计"""
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=503, detail="PPT 功能不可用")

    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="只支持 .pptx 文件")

    try:
        content = await file.read()
        prs = Presentation(io.BytesIO(content))

        slides_data = []
        for slide_idx, slide in enumerate(prs.slides):
            bg_color = None
            try:
                if slide.background.fill.type is not None and slide.background.fill.type == 1:
                    bg_color = f"#{slide.background.fill.fore_color.rgb}"
            except:
                pass

            shapes = []
            for shape in slide.shapes:
                shapes.extend(_extract_shapes_recursive(shape))

            slides_data.append(SlideData(
                index=slide_idx + 1,
                shapes=[SlideShape(**s) for s in shapes],
                background=bg_color,
            ))

        theme_name = _detect_theme_from_pptx(prs)
        theme = DesignPresets.get(theme_name)

        preview = PPTPreviewData(
            filename=file.filename,
            total_slides=len(prs.slides),
            slide_width=round(prs.slide_width.emu / 914400 * 96, 2) if hasattr(prs, 'slide_width') else 960,
            slide_height=round(prs.slide_height.emu / 914400 * 96, 2) if hasattr(prs, 'slide_height') else 540,
            slides=slides_data,
            design_system=BrandStyle(theme=theme),
        )

        return _serialize_preview(preview)

    except Exception as e:
        logger.error(f"PPT 提取失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/preview/html")
async def generate_ppt_html_preview(file: UploadFile = File(...)):
    """生成 PPT 的 HTML 预览（保留原始样式）"""
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=503, detail="PPT 功能不可用")

    try:
        content = await file.read()
        prs = Presentation(io.BytesIO(content))

        slide_width_px = round(prs.slide_width.emu / 914400 * 96) if hasattr(prs, 'slide_width') else 960
        slide_height_px = round(prs.slide_height.emu / 914400 * 96) if hasattr(prs, 'slide_height') else 540
        scale = min(800 / slide_width_px, 600 / slide_height_px)
        display_width = int(slide_width_px * scale)
        display_height = int(slide_height_px * scale)

        slides_html = []
        for slide_idx, slide in enumerate(prs.slides):
            bg_color = "#ffffff"
            try:
                if slide.background.fill.type is not None and slide.background.fill.type == 1:
                    bg_color = f"#{slide.background.fill.fore_color.rgb}"
            except:
                pass

            shapes_html = []
            for shape in slide.shapes:
                html = _gen_shape_html(shape, scale)
                if html:
                    shapes_html.append(html)

            slides_html.append(f"""
                <div class="slide" style="width:{display_width}px;height:{display_height}px;">
                    <div class="slide-num">{slide_idx + 1}</div>
                    <div class="slide-content" style="width:100%;height:100%;background:{bg_color};position:relative;">
                        {''.join(shapes_html)}
                    </div>
                </div>""")

        html = f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI','Microsoft YaHei',sans-serif;background:#1a1a2e;padding:20px}}
.ppt-container{{max-width:900px;margin:0 auto}}
.slide{{background:white;margin-bottom:20px;border-radius:8px;overflow:hidden;box-shadow:0 4px 20px rgba(0,0,0,.3);position:relative}}
.slide-num{{position:absolute;top:10px;right:10px;background:rgba(0,0,0,.5);color:white;padding:4px 12px;border-radius:12px;font-size:12px;z-index:10}}
.shape-text{{word-wrap:break-word;line-height:1.4}}
.shape-table{{border-collapse:collapse;width:100%;height:100%}}
.shape-table td{{border:1px solid #ccc;padding:4px;font-size:12px}}
</style></head><body><div class="ppt-container">{''.join(slides_html)}</div></body></html>"""

        return {"html": html, "total_slides": len(prs.slides)}

    except Exception as e:
        logger.error(f"HTML 预览生成失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


def _gen_shape_html(shape, scale: float) -> Optional[str]:
    try:
        left = int(shape.left.emu / 914400 * 96 * scale) if hasattr(shape, 'left') else 0
        top = int(shape.top.emu / 914400 * 96 * scale) if hasattr(shape, 'top') else 0
        width = int(shape.width.emu / 914400 * 96 * scale) if hasattr(shape, 'width') else 100
        height = int(shape.height.emu / 914400 * 96 * scale) if hasattr(shape, 'height') else 50

        fill_style = ""
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None and shape.fill.type == 1:
                fill_style = f"background:#{shape.fill.fore_color.rgb};"
        except:
            pass

        content = ""
        text_style = ""

        if shape.has_text_frame:
            texts = []
            for para in shape.text_frame.paragraphs:
                para_text = para.text
                if para_text:
                    texts.append(para_text)

            if texts:
                content = "<br>".join(texts)
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                text_style += f"font-size:{int(run.font.size.pt * scale)}px;"
                            if run.font.bold:
                                text_style += "font-weight:bold;"
                            if run.font.color and run.font.color.rgb:
                                text_style += f"color:#{run.font.color.rgb};"
                            break
                        break
                except:
                    text_style += f"font-size:{int(18 * scale)}px;color:#333;"

        if shape.has_table:
            rows_html = []
            for row in shape.table.rows:
                cells_html = [f"<td>{cell.text}</td>" for cell in row.cells]
                rows_html.append(f"<tr>{''.join(cells_html)}</tr>")
            content = f"<table class='shape-table'>{''.join(rows_html)}</table>"

        if content:
            return f"""<div class="shape" style="position:absolute;left:{left}px;top:{top}px;width:{width}px;height:{height}px;{fill_style}">
                <div class="shape-text" style="{text_style}width:100%;height:100%;overflow:hidden;">{content}</div>
            </div>"""
        return None
    except Exception as e:
        logger.warning(f"形状 HTML 生成失败: {e}")
        return None


@router.post("/reconstruct")
async def reconstruct_ppt(data: dict = Body(...)):
    """
    根据编辑后的幻灯片数据重新生成 PPTX
    保留原有设计系统样式
    """
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=503, detail="PPT 功能不可用")

    try:
        from pptx import Presentation as Pres
        from pptx.util import Inches, Pt

        slides_data = data.get("slides", [])
        theme_name = data.get("theme", "professional")
        filename = data.get("filename", "reconstructed.pptx")

        theme = DesignPresets.get(theme_name)
        prs = Pres()
        prs.slide_width = Emu(int(data.get("slide_width", 960) / 96 * 914400))
        prs.slide_height = Emu(int(data.get("slide_height", 540) / 96 * 914400))

        for s in slides_data:
            layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(layout)

            if s.get("background"):
                try:
                    from pptx.oxml.ns import qn
                    bg = slide.background
                    fill_elem = bg._element.find(qn('p:bgPr'))
                    if fill_elem is None:
                        bg_elem = bg._element.makeelement(qn('p:bgPr'), {})
                        bg._element.insert(0, bg_elem)
                        fill_elem = bg_elem
                    from pptx.oxml.ns import qn as qn2
                    from lxml import etree
                    solid = fill_elem.find(qn2('a:solidFill'))
                    if solid is None:
                        solid = etree.SubElement(fill_elem, qn2('a:solidFill'))
                    srgb = solid.find(qn2('a:srgbClr'))
                    if srgb is None:
                        srgb = etree.SubElement(solid, qn2('a:srgbClr'))
                    h = s["background"].lstrip("#")
                    srgb.set('val', h)
                except Exception as bg_err:
                    logger.warning(f"设置背景失败: {bg_err}")

            for shape_data in s.get("shapes", []):
                try:
                    left = Inches(shape_data.get("left", 0) / 96)
                    top = Inches(shape_data.get("top", 0) / 96)
                    width = Inches(shape_data.get("width", 200) / 96)
                    height = Inches(shape_data.get("height", 50) / 96)

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True

                    paragraphs = shape_data.get("paragraphs", [{"text": shape_data.get("text", "")}])
                    for pi, para_data in enumerate(paragraphs):
                        if pi == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()

                        runs = para_data.get("runs", [{"text": para_data.get("text", "")}])
                        for ri, run_data in enumerate(runs):
                            if ri == 0:
                                run = p.runs[0] if p.runs else p.add_run()
                            else:
                                run = p.add_run()
                            run.text = run_data.get("text", "")

                            try:
                                fs = run_data.get("font_size")
                                if fs:
                                    run.font.size = Pt(fs)
                                if run_data.get("bold"):
                                    run.font.bold = True
                                if run_data.get("italic"):
                                    run.font.italic = True
                                fc = run_data.get("color")
                                if fc:
                                    h = fc.lstrip("#")
                                    run.font.color.rgb = RGBColor(
                                        int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                                    )
                            except:
                                pass
                except Exception as shape_err:
                    logger.warning(f"重建形状失败: {shape_err}")

        output_dir = Path("C:/D/zhiyi/generated")
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / filename
        prs.save(str(output_path))

        return {"success": True, "filename": filename}

    except Exception as e:
        logger.error(f"PPT 重建失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


def _serialize_preview(preview: PPTPreviewData) -> dict:
    """序列化 PPTPreviewData 为字典"""
    return {
        "filename": preview.filename,
        "total_slides": preview.total_slides,
        "slide_width": preview.slide_width,
        "slide_height": preview.slide_height,
        "slides": [
            {
                "index": s.index,
                "background": s.background,
                "shapes": [
                    {
                        "type": sh.type,
                        "left": sh.left,
                        "top": sh.top,
                        "width": sh.width,
                        "height": sh.height,
                        "text": sh.text,
                        "font_size": sh.font_size,
                        "font_bold": sh.font_bold,
                        "font_color": sh.font_color,
                        "fill_color": sh.fill_color,
                        "table": sh.table,
                        "image_url": sh.image_url,
                        "paragraphs": getattr(sh, "paragraphs", None) or (
                            [{"text": sh.text}] if sh.text else []
                        ),
                        "default_font_size": getattr(sh, "default_font_size", None),
                        "default_font_color": getattr(sh, "default_font_color", None),
                        "default_font_name": getattr(sh, "default_font_name", None),
                    }
                    for sh in s.shapes
                ],
                "notes": s.notes,
            }
            for s in preview.slides
        ],
        "design_system": preview.design_system.to_dict() if preview.design_system else None,
    }
