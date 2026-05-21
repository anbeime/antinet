"""
知识管理路由
提供知识库的 CRUD 接口
"""
from fastapi import APIRouter, HTTPException, UploadFile, File, Request
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any
from datetime import datetime
import logging
import tempfile
import os
import re
import json
import html
import hashlib
from pathlib import Path
import time
import asyncio
from functools import lru_cache

from config import settings
from database import DatabaseManager

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/knowledge", tags=["knowledge", "知识网络"])


def sanitize_html(content: str) -> str:
    """过滤 HTML 内容防止 XSS 攻击"""
    if not content:
        return content
    # 解码 HTML 实体以处理输入
    text = content
    # 允许的安全标签和属性
    allowed_tags = ['b', 'i', 'u', 'strong', 'em', 'code', 'pre', 'br', 'p', 'ul', 'ol', 'li', 'a', 'span']
    allowed_attrs = {'a': ['href', 'title'], 'span': ['class']}
    
    # 简单处理：移除危险属性
    dangerous_patterns = [
        (r'on\w+\s*=', ''),  # 移除事件处理器
        (r'javascript:', ''),  # 移除 JS 协议
        (r'data:', ''),       # 移除 data 协议
    ]
    for pattern, repl in dangerous_patterns:
        text = re.sub(pattern, repl, text, flags=re.IGNORECASE)
    
    return text


# ==================== 知识网络入口 ====================

class NetworkCardsRequest(BaseModel):
    """知识网络卡片请求"""
    topic: str = Field(..., description="主题/查询词")
    color_filter: Optional[str] = Field(default=None, description="颜色过滤: blue/green/yellow/red")
    mode: str = Field(default="auto", description="模式: auto(自动生成)|manual(手动选择)")
    limit: int = Field(default=10000, description="返回数量")


class NetworkSuggestion(BaseModel):
    """网络建议"""
    card_id: str
    card_type: str
    title: str
    content: str
    category: str
    reason: str  # 为什么推荐这张卡片


class NetworkGenerateRequest(BaseModel):
    """网络生成请求"""
    topic: str
    card_ids: Optional[List[str]] = Field(default=None, description="手动选择的卡片ID")
    auto_generate: bool = Field(default=True, description="是否自动生成相关卡片")
    target_type: str = Field(default="both", description="目标: kg(图谱)|mindmap(导图)|both")


@router.post("/network/cards")
async def get_network_cards(request: NetworkCardsRequest):
    """
    知识网络 - 获取相关卡片
    根据主题查找或生成相关卡片
    """
    conn = db_manager.get_connection()
    cursor = conn.cursor()
    
    try:
        if request.color_filter:
            cursor.execute("""
                SELECT id, card_type, title, content, category, created_at
                FROM knowledge_cards
                WHERE (title LIKE ? OR content LIKE ?) 
                AND card_type = ?
                ORDER BY created_at DESC
                LIMIT ?
            """, [f"%{request.topic}%", f"%{request.topic}%", request.color_filter, request.limit])
        else:
            cursor.execute("""
                SELECT id, card_type, title, content, category, created_at
                FROM knowledge_cards
                WHERE title LIKE ? OR content LIKE ?
                ORDER BY created_at DESC
                LIMIT ?
            """, [f"%{request.topic}%", f"%{request.topic}%", request.limit])
        
        rows = cursor.fetchall()
        cards = []
        for row in rows:
            cards.append({
                "card_id": str(row["id"]),
                "card_type": row["card_type"],
                "title": row["title"],
                "content": row["content"],
                "category": row["category"],
                "created_at": row["created_at"],
                "reason": f"与主题「{request.topic}」相关"
            })
        
        return {
            "topic": request.topic,
            "cards": cards,
            "total": len(cards),
            "mode": request.mode
        }
    finally:
        conn.close()


@router.get("/network/suggest")
@router.post("/network/suggest")
async def suggest_network_cards(topic: str, limit: int = 10000):
    """
    知识网络 - AI联想推荐卡片
    基于主题智能推荐相关卡片（语义搜索）
    """
    try:
        from routes import vector_search
        vector_search.set_db_manager(db_manager)
        
        # 使用混合搜索：向量 + 关键词
        results = vector_search.search_hybrid(topic, limit=min(limit, 20))
        
        if not results:
            # 回退到关键词搜索
            results = vector_search.fallback_keyword_search(topic, limit=min(limit, 20))
        
        return {
            "topic": topic,
            "total": len(results),
            "suggestions": [
                {
                    "card_id": r.id,
                    "title": r.title,
                    "content": r.content[:200] + "..." if len(r.content) > 200 else r.content,
                    "card_type": r.card_type,
                    "score": r.score
                }
                for r in results
            ]
        }
    except Exception as e:
        logger.error(f"智能推荐失败: {e}")
        return {"topic": topic, "suggestions": [], "error": str(e)}


@router.post("/network/generate")
async def generate_network(request: NetworkGenerateRequest):
    """
    知识网络 - 生成知识网络
    将选中的卡片（手动或自动）生成知识图谱或思维导图
    """
    conn = db_manager.get_connection()
    cursor = conn.cursor()
    
    try:
        created_entities = []
        created_relations = []
        
        # 获取选中的卡片或自动搜索
        if request.card_ids:
            placeholders = ",".join(["?" for _ in request.card_ids])
            cursor.execute(f"""
                SELECT id, card_type, title, content
                FROM knowledge_cards
                WHERE id IN ({placeholders})
            """, request.card_ids)
        else:
            cursor.execute("""
                SELECT id, card_type, title, content
                FROM knowledge_cards
                WHERE title LIKE ? OR content LIKE ?
                ORDER BY RANDOM()
                LIMIT 12
            """, [f"%{request.topic}%", f"%{request.topic}%"])
        
        cards = cursor.fetchall()
        
        # 1. 创建图谱实体
        if request.target_type in ("kg", "both"):
            for card in cards[:12]:
                card_id = card['id']
                cursor.execute("""
                    INSERT OR IGNORE INTO kg_entities (entity_id, name, entity_type, description)
                    VALUES (?, ?, ?, ?)
                """, [f"card_{card_id}", card["title"][:50], card["card_type"], card["content"][:200]])
                created_entities.append(f"card_{card_id}")
            
            # 创建主题实体
            cursor.execute("""
                INSERT OR IGNORE INTO kg_entities (entity_id, name, entity_type, description)
                VALUES (?, ?, ?, ?)
            """, [f"topic_{request.topic[:20]}", request.topic[:50], "主题", request.topic])
            created_entities.append(f"topic_{request.topic[:20]}")
            
            # 建立关系
            for entity_id in created_entities[1:]:
                cursor.execute("""
                    INSERT OR IGNORE INTO kg_relations (relation_id, source_id, target_id, relation_type)
                    VALUES (?, ?, ?, ?)
                """, [f"rel_{created_entities[0][-10:]}_{entity_id[-10:]}", created_entities[0], entity_id, "关联"])
                created_relations.append(f"rel_{created_entities[0][-10:]}_{entity_id[-10:]}")
        
        # 2. 创建思维导图
        if request.target_type in ("mindmap", "both"):
            root_node = {
                "id": f"root_{int(datetime.now().timestamp())}",
                "text": request.topic,
                "children": [],
                "color": "#8b5cf6"
            }
            
            for i, card in enumerate(cards[:12]):
                color_map = {"blue": "#3b82f6", "green": "#22c55e", "yellow": "#eab308", "red": "#ef4444"}
                root_node["children"].append({
                    "id": f"node_{i}",
                    "text": card["title"][:30],
                    "children": [],
                    "color": color_map.get(card["card_type"], "#8b5cf6")
                })
            
            cursor.execute("""
                INSERT INTO mindmaps (name, root_node, created_at)
                VALUES (?, ?, datetime('now'))
            """, [f"知识网络-{request.topic[:20]}", json.dumps(root_node)])
            mindmap_id = cursor.lastrowid
        
        conn.commit()
        
        return {
            "status": "generated",
            "topic": request.topic,
            "entities_created": len(created_entities),
            "relations_created": len(created_relations),
            "mindmap_id": mindmap_id if request.target_type in ("mindmap", "both") else None
        }
    except Exception as e:
        logger.error(f"[KnowledgeNetwork] 生成失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        conn.close()


# 创建数据库管理器实例
db_manager = DatabaseManager(settings.DB_PATH)


class SearchRequest(BaseModel):
    """知识库搜索请求"""
    keyword: str = Field(..., description="搜索关键词")
    limit: int = Field(10000, description="返回数量限制")


@router.get("/graph")
async def get_knowledge_graph(
    card_type: Optional[str] = None,
    project_id: Optional[int] = None,
    limit: int = 500
):
    """
    获取知识图谱数据（带缓存）
    """
    from functools import partial
    
    # 初始化缓存
    if not hasattr(get_knowledge_graph, '_cache'):
        get_knowledge_graph._cache = {}
        get_knowledge_graph._cache_time = {}
        get_knowledge_graph._clear = lambda: (
            get_knowledge_graph._cache.clear(),
            get_knowledge_graph._cache_time.clear()
        )
    
    # 检查缓存（5分钟内有效）
    cache_key = f"graph_{card_type}_{project_id}_{limit}"
    current_time = time.time()
    
    if cache_key in get_knowledge_graph._cache:
        cached_time = get_knowledge_graph._cache_time.get(cache_key, 0)
        if current_time - cached_time < 300:
            logger.info(f"[GRAPH] 使用缓存: {cache_key}")
            return get_knowledge_graph._cache[cache_key]
    
    logger.info(f"[GRAPH] 构建新图谱: {cache_key}")
    try:
        from services.skill_system import get_skill_registry
        
        # 获取技能注册表
        registry = get_skill_registry()
        
        # 获取卡片
        conn = db_manager.get_connection()
        cursor = conn.cursor()
        
        query = "SELECT * FROM knowledge_cards WHERE 1=1"
        params = []
        
        if card_type:
            query += " AND card_type = ?"
            params.append(card_type)
        
        if project_id:
            query += " AND project_id = ?"
            params.append(project_id)
        
        # 添加排序和限制
        query += " ORDER BY created_at DESC LIMIT ?"
        params.append(limit)
        
        cursor.execute(query, params)
        rows = cursor.fetchall()
        conn.close()
        
        # 转换为字典列表
        cards = [dict(row) for row in rows]
        
        # 调用知识图谱可视化技能
        result = await registry.execute_skill(
            "knowledge_graph_visualization",
            cards=cards
        )
        
        graph_data = result.get("result", {})
        
        # 缓存结果
        get_knowledge_graph._cache[cache_key] = graph_data
        get_knowledge_graph._cache_time[cache_key] = current_time
        
        return graph_data
        
    except Exception as e:
        logger.error(f"获取知识图谱失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


class KnowledgeCard(BaseModel):
    """知识卡片模型"""
    model_config = {"populate_by_name": True}
    
    id: Optional[int] = None
    type: str = Field(validation_alias='card_type')  # 使用 type 内部，card_type 来自前端
    title: Optional[str] = None  # 允许为空，由系统自动生成
    content: str
    source: Optional[str] = None
    url: Optional[str] = None
    category: Optional[str] = None
    project_id: Optional[int] = None  # 关联的专题ID
    related_cards: Optional[List[int]] = []  # 关联的卡片ID列表
    address: Optional[str] = None  # Antinet 地址
    source_file_id: Optional[str] = None  # 源文件ID（文件导入时关联）
    images: Optional[List[Dict[str, Any]]] = []  # 图片列表 [{id, filename, path, url, size}]
    
    def get_valid_category(self) -> str:
        """获取有效的 category"""
        valid_categories = {'事实', '解释', '风险', '行动'}
        if self.category and self.category in valid_categories:
            return self.category
        # 根据卡片类型推断 category
        type_to_category = {
            'blue': '事实',
            'green': '解释',
            'yellow': '风险',
            'red': '行动'
        }
        return type_to_category.get(self.type, '事实')
    
    def get_or_generate_title(self) -> str:
        """获取或自动生成标题"""
        if self.title and self.title.strip():
            return self.title.strip()
        
        # 从内容中提取标题
        if self.content:
            # 取内容的前30个字符作为标题
            preview = self.content.strip()[:30]
            # 去掉可能的标点符号
            preview = preview.rstrip('。！?？,，')
            if preview:
                return preview + "..."
        
        # 默认标题
        type_to_title = {
            'blue': '新事实',
            'green': '新解释',
            'yellow': '新风险',
            'red': '新行动'
        }
        return type_to_title.get(self.type, '新卡片')


class KnowledgeSource(BaseModel):
    """知识来源模型"""
    id: Optional[int] = None
    source_path: str
    source_type: str
    total_cards: int = 0


@router.get("/cards")
async def get_cards(
    card_type: Optional[str] = None,
    category: Optional[str] = None,
    limit: int = 10000,
    offset: int = 0
):
    """
    获取知识卡片列表

    Args:
        card_type: 卡片类型过滤（blue/green/yellow/red）
        category: 分类过滤
        limit: 返回数量限制
        offset: 偏移量

    Returns:
        卡片列表
    """
    conn = db_manager.get_connection()
    cursor = conn.cursor()

    query = "SELECT * FROM knowledge_cards WHERE 1=1"
    params = []

    if card_type:
        query += " AND card_type = ?"
        params.append(card_type)

    if category:
        query += " AND category = ?"
        params.append(category)

    query += " ORDER BY created_at DESC LIMIT ? OFFSET ?"
    params.extend([limit, offset])

    cursor.execute(query, params)
    columns = [description[0] for description in cursor.description]
    cards = []
    for row in cursor.fetchall():
        card_dict = dict(zip(columns, row))
        # 解析 related_cards JSON
        if card_dict.get('related_cards'):
            try:
                card_dict['related_cards'] = json.loads(card_dict['related_cards'])
            except:
                card_dict['related_cards'] = []
        else:
            card_dict['related_cards'] = []
        # 解析 images JSON
        if card_dict.get('images'):
            try:
                card_dict['images'] = json.loads(card_dict['images'])
            except:
                card_dict['images'] = []
        else:
            card_dict['images'] = []
        cards.append(card_dict)

    conn.close()
    return {"cards": cards, "total": len(cards)}


@router.get("/cards/{card_id}")
async def get_card(card_id: int):
    """
    获取单个知识卡片

    Args:
        card_id: 卡片ID

    Returns:
        卡片详情
    """
    conn = db_manager.get_connection()
    cursor = conn.cursor()

    cursor.execute("SELECT * FROM knowledge_cards WHERE id = ?", (card_id,))
    card = cursor.fetchone()

    conn.close()

    if not card:
        raise HTTPException(status_code=404, detail="卡片不存在")

    card_dict = dict(card)
    # 解析 related_cards JSON
    if card_dict.get('related_cards'):
        try:
            card_dict['related_cards'] = json.loads(card_dict['related_cards'])
        except:
            card_dict['related_cards'] = []
    else:
        card_dict['related_cards'] = []
    # 解析 images JSON
    if card_dict.get('images'):
        try:
            card_dict['images'] = json.loads(card_dict['images'])
        except:
            card_dict['images'] = []
    else:
        card_dict['images'] = []

    return card_dict


@router.post("/cards")
async def create_card(card: KnowledgeCard):
    """
    创建新的知识卡片

    Args:
        card: 卡片数据

    Returns:
        创建的卡片
    """
    logger.info(f"[CREATE_CARD] 收到创建卡片请求: {card.dict()}")

    conn = db_manager.get_connection()
    cursor = conn.cursor()

    try:
        logger.info(f"[CREATE_CARD] 准备插入数据库，type={card.type}")

        # 使用有效的 category
        valid_category = card.get_valid_category()
        
        # 自动生成标题（如果未提供）
        card_title = card.get_or_generate_title()
        
        # 自动生成地址（如果未提供）
        card_address = card.address
        if not card_address:
            type_prefix = {'blue': 'A', 'green': 'B', 'yellow': 'C', 'red': 'D'}
            prefix = type_prefix.get(card.type, 'X')
            cursor.execute("SELECT COUNT(*) FROM knowledge_cards WHERE card_type = ?", (card.type,))
            count = cursor.fetchone()[0] + 1
            card_address = f"{prefix}{count}"
        
        # 使用正确的字段名 card_type（与数据库表结构一致）
        related_cards_json = json.dumps(card.related_cards) if card.related_cards else None
        images_json = json.dumps(card.images) if card.images else '[]'  # 图片列表JSON
        
        # 过滤 HTML 内容防止 XSS
        safe_content = sanitize_html(card.content) if card.content else ''
        
        cursor.execute('''
            INSERT INTO knowledge_cards (card_type, title, content, category, project_id, related_cards, address, images)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        ''', (
            card.type,
            card_title,
            safe_content,
            valid_category,
            card.project_id,
            related_cards_json,
            card_address,
            images_json
        ))

        new_card_id = cursor.lastrowid

        # 同步双向链接到 card_backlinks 表
        if card.related_cards:
            try:
                for target_id in card.related_cards:
                    cursor.execute("""
                        INSERT OR IGNORE INTO card_backlinks (source_card_id, target_card_id, link_text)
                        VALUES (?, ?, ?)
                    """, (new_card_id, target_id, 'manual'))
                    # 双向
                    cursor.execute("""
                        INSERT OR IGNORE INTO card_backlinks (source_card_id, target_card_id, link_text)
                        VALUES (?, ?, ?)
                    """, (target_id, new_card_id, 'manual'))
            except Exception as e:
                logger.warning(f"同步backlinks失败（非致命）: {e}")

        conn.commit()

        logger.info(f"[CREATE_CARD] 插入成功，lastrowid={new_card_id}")

        # 获取新插入的卡片
        cursor.execute("SELECT * FROM knowledge_cards WHERE id = ?", (cursor.lastrowid,))
        new_card = dict(cursor.fetchone())

        conn.close()
        logger.info(f"[CREATE_CARD] 返回新卡片: {new_card}")
        
        # 清除图谱缓存
        if hasattr(get_knowledge_graph, '_clear'):
            get_knowledge_graph._clear()
        
        return new_card

    except Exception as e:
        logger.error(f"[CREATE_CARD] 创建失败: {e}", exc_info=True)
        conn.close()
        raise HTTPException(status_code=400, detail=f"创建失败: {str(e)}")


@router.put("/cards/{card_id}")
async def update_card(card_id: int, card: KnowledgeCard):
    """
    更新知识卡片

    Args:
        card_id: 卡片ID
        card: 更新的卡片数据

    Returns:
        更新后的卡片
    """
    logger.info(f"[UPDATE_CARD] 收到更新卡片请求: id={card_id}, data={card.dict()}")

    conn = db_manager.get_connection()
    cursor = conn.cursor()

    try:
        # 检查卡片是否存在
        cursor.execute("SELECT * FROM knowledge_cards WHERE id = ?", (card_id,))
        existing_card = cursor.fetchone()
        
        if not existing_card:
            conn.close()
            raise HTTPException(status_code=404, detail="卡片不存在")

        # 更新卡片 — 保留未被显式传入的字段（如 project_id）
        related_cards_json = json.dumps(card.related_cards) if card.related_cards else None
        
        # 过滤 HTML 内容防止 XSS
        safe_content = sanitize_html(card.content) if card.content else None
        
        # 如果 project_id 未传入（None），保留数据库中的现有值
        existing_project_id = dict(existing_card).get('project_id') if existing_card else None
        final_project_id = card.project_id if card.project_id is not None else existing_project_id
        # 如果 address 未传入，保留现有值
        existing_address = dict(existing_card).get('address') if existing_card else None
        final_address = card.address if card.address else existing_address
        
        cursor.execute('''
            UPDATE knowledge_cards 
            SET card_type = ?, title = ?, content = ?, category = ?, project_id = ?, related_cards = ?, address = ?
            WHERE id = ?
        ''', (
            card.type,
            card.title,
            safe_content if safe_content is not None else dict(existing_card).get('content'),
            card.category,
            final_project_id,
            related_cards_json,
            final_address,
            card_id
        ))

        # 同步双向链接到 card_backlinks 表
        if card.related_cards is not None:
            try:
                # 获取旧关联（从 backlinks 表）
                cursor.execute("""
                    SELECT target_card_id FROM card_backlinks WHERE source_card_id = ? AND link_text = 'manual'
                """, (card_id,))
                old_related = set(row[0] for row in cursor.fetchall())
                
                new_related = set(card.related_cards)
                
                # 新增关联
                for target_id in (new_related - old_related):
                    cursor.execute("""
                        INSERT OR IGNORE INTO card_backlinks (source_card_id, target_card_id, link_text)
                        VALUES (?, ?, ?)
                    """, (card_id, target_id, 'manual'))
                    # 双向：目标卡片也链接回来
                    cursor.execute("""
                        INSERT OR IGNORE INTO card_backlinks (source_card_id, target_card_id, link_text)
                        VALUES (?, ?, ?)
                    """, (target_id, card_id, 'manual'))
                
                # 移除关联
                for target_id in (old_related - new_related):
                    cursor.execute("""
                        DELETE FROM card_backlinks 
                        WHERE ((source_card_id = ? AND target_card_id = ?) 
                           OR (source_card_id = ? AND target_card_id = ?))
                           AND link_text = 'manual'
                    """, (card_id, target_id, target_id, card_id))
            except Exception as e:
                logger.warning(f"同步backlinks失败（非致命）: {e}")

        conn.commit()

        # 获取更新后的卡片
        cursor.execute("SELECT * FROM knowledge_cards WHERE id = ?", (card_id,))
        updated_card = dict(cursor.fetchone())

        conn.close()
        logger.info(f"[UPDATE_CARD] 更新成功: {updated_card}")
        
        # 清除图谱缓存
        if hasattr(get_knowledge_graph, '_clear'):
            get_knowledge_graph._clear()
        
        return updated_card

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[UPDATE_CARD] 更新失败: {e}", exc_info=True)
        conn.close()
        raise HTTPException(status_code=400, detail=f"更新失败: {str(e)}")


@router.delete("/cards/{card_id}")
async def delete_card(card_id: int):
    """
    删除知识卡片

    Args:
        card_id: 卡片ID

    Returns:
        删除结果
    """
    conn = db_manager.get_connection()
    cursor = conn.cursor()

    cursor.execute("DELETE FROM knowledge_cards WHERE id = ?", (card_id,))
    conn.commit()

    conn.close()
    
    # 清除图谱缓存
    if hasattr(get_knowledge_graph, '_clear'):
        get_knowledge_graph._clear()

    return {"success": True, "message": "卡片已删除"}


@router.get("/stats")
async def get_stats():
    """
    获取知识库统计信息

    Returns:
        统计信息
    """
    try:
        conn = db_manager.get_connection()
        cursor = conn.cursor()

        # 总卡片数
        cursor.execute("SELECT COUNT(*) FROM knowledge_cards")
        total_cards = cursor.fetchone()[0]

        # 按类型分组 - 使用 card_type 字段（数据库字段名）
        cursor.execute("SELECT card_type, COUNT(*) as count FROM knowledge_cards WHERE card_type IS NOT NULL GROUP BY card_type")
        cards_by_type = {row[0]: row[1] for row in cursor.fetchall() if row[0] is not None}

        # 按分类分组
        cursor.execute("SELECT category, COUNT(*) as count FROM knowledge_cards WHERE category IS NOT NULL GROUP BY category")
        cards_by_category = {row[0]: row[1] for row in cursor.fetchall() if row[0] is not None}

        conn.close()

        return {
            "total_cards": total_cards,
            "cards_by_type": cards_by_type,
            "cards_by_category": cards_by_category
        }
    except Exception as e:
        logger.error(f"获取统计信息失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/search")
async def search_cards(request: SearchRequest):
    """
    搜索知识卡片

    Args:
        request: 搜索请求（包含关键词和限制）

    Returns:
匹配的卡片列表
    """
    try:
        conn = db_manager.get_connection()
        cursor = conn.cursor()

        keyword = request.keyword
        limit = request.limit

        cursor.execute('''
            SELECT * FROM knowledge_cards
            WHERE title LIKE ? OR content LIKE ?
            ORDER BY created_at DESC
            LIMIT ?
        ''', (f'%{keyword}%', f'%{keyword}%', limit))

        cards = [dict(row) for row in cursor.fetchall()]
        conn.close()

        return cards
    except Exception as e:
        logger.error(f"搜索失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"搜索失败: {str(e)}")


# ==================== VCP TagMemo 风格标签系统 ====================

class TagUpdateRequest(BaseModel):
    """标签更新请求"""
    card_id: int
    tags: Optional[List[str]] = None
    core_tags: Optional[List[str]] = None
    tag_weights: Optional[Dict[str, float]] = None


class RecallRequest(BaseModel):
    """智能召回请求"""
    query: str
    memory_type: str = "light"  
    core_tags_boost: bool = True
    limit: int = 10


@router.get("/tags")
async def get_all_tags():
    """获取所有标签（含权重统计）- TagMemo 风格"""
    conn = db_manager.get_connection()
    cursor = conn.cursor()
    
    all_tags = {}
    all_core_tags = {}
    
    cursor.execute("SELECT id, tags, core_tags, tag_weights FROM knowledge_cards")
    for row in cursor.fetchall():
        card_id, tags_json, core_tags_json, weights_json = row
        
        try:
            tags = json.loads(tags_json) if tags_json else []
            for tag in tags:
                all_tags[tag] = all_tags.get(tag, 0) + 1
        except:
            pass
        
        try:
            core_tags = json.loads(core_tags_json) if core_tags_json else []
            for tag in core_tags:
                all_core_tags[tag] = all_core_tags.get(tag, 0) + 1
        except:
            pass
    
    conn.close()
    
    return {
        "tags": [{"name": k, "count": v} for k, v in sorted(all_tags.items(), key=lambda x: -x[1])],
        "core_tags": [{"name": k, "count": v} for k, v in sorted(all_core_tags.items(), key=lambda x: -x[1])]
    }


@router.post("/cards/{card_id}/tags")
async def update_card_tags(card_id: int, request: TagUpdateRequest):
    """更新卡片标签 - TagMemo 风格"""
    conn = db_manager.get_connection()
    cursor = conn.cursor()
    
    cursor.execute("SELECT id FROM knowledge_cards WHERE id = ?", (card_id,))
    if not cursor.fetchone():
        conn.close()
        raise HTTPException(status_code=404, detail="卡片不存在")
    
    updates = []
    params = []
    
    if request.tags is not None:
        updates.append("tags = ?")
        params.append(json.dumps(request.tags))
    
    if request.core_tags is not None:
        updates.append("core_tags = ?")
        params.append(json.dumps(request.core_tags))
    
    if request.tag_weights is not None:
        updates.append("tag_weights = ?")
        params.append(json.dumps(request.tag_weights))
    
    if updates:
        updates.append("updated_at = ?")
        params.append(datetime.now().isoformat())
        params.append(card_id)
        
        cursor.execute(f"UPDATE knowledge_cards SET {', '.join(updates)} WHERE id = ?", params)
        conn.commit()
    
    conn.close()
    return {"success": True, "message": "标签已更新"}


@router.post("/recall")
async def recall_cards(request: RecallRequest):
    """智能召回 - TagMemo LIF-Router 风格"""
    conn = db_manager.get_connection()
    cursor = conn.cursor()
    
    query = request.query.lower()
    cards = []
    
    if request.memory_type == "light":
        cursor.execute('''
            SELECT id, title, content, tags, core_tags, tag_weights, coherence_score, access_count
            FROM knowledge_cards
            WHERE title LIKE ? OR content LIKE ? OR tags LIKE ?
            ORDER BY CASE WHEN core_tags LIKE ? THEN 0 ELSE 1 END, access_count DESC
            LIMIT ?
        ''', (f'%{query}%', f'%{query}%', f'%{query}%', f'%{query}%', request.limit))
    elif request.memory_type == "deep":
        cursor.execute('''
            SELECT id, title, content, tags, core_tags, tag_weights, coherence_score, access_count
            FROM knowledge_cards
            WHERE title LIKE ? OR content LIKE ? OR tags LIKE ? OR core_tags LIKE ?
            ORDER BY coherence_score DESC, access_count DESC
            LIMIT ?
        ''', (f'%{query}%', f'%{query}%', f'%{query}%', f'%{query}%', request.limit))
    else:
        cursor.execute('''
            SELECT id, title, content, tags, core_tags, tag_weights, coherence_score, access_count
            FROM knowledge_cards
            WHERE (title LIKE ? OR content LIKE ?) AND (tags LIKE ? OR core_tags LIKE ?)
            ORDER BY (coherence_score * 0.5 + access_count * 0.1) DESC
            LIMIT ?
        ''', (f'%{query}%', f'%{query}%', f'%{query}%', f'%{query}%', request.limit))
    
    for row in cursor.fetchall():
        cards.append({
            "id": row[0], "title": row[1], "content": row[2],
            "tags": json.loads(row[3]) if row[3] else [],
            "core_tags": json.loads(row[4]) if row[4] else [],
            "tag_weights": json.loads(row[5]) if row[5] else {},
            "coherence_score": row[6] or 0.0, "access_count": row[7] or 0
        })
    
    conn.close()
    return {"memory_type": request.memory_type, "query": request.query, "results": cards, "count": len(cards)}


@router.post("/cards/{card_id}/access")
async def track_card_access(card_id: int):
    """追踪卡片访问"""
    conn = db_manager.get_connection()
    cursor = conn.cursor()
    
    cursor.execute("""
        UPDATE knowledge_cards 
        SET access_count = COALESCE(access_count, 0) + 1, last_accessed = ?
        WHERE id = ?
    """, (datetime.now().isoformat(), card_id))
    
    conn.commit()
    conn.close()
    return {"success": True}


@router.get("/sources")
async def get_sources():
    """
    获取知识来源列表

    Returns:
        来源列表
    """
    conn = db_manager.get_connection()
    cursor = conn.cursor()

    cursor.execute("SELECT * FROM knowledge_sources ORDER BY last_imported DESC")
    sources = [dict(row) for row in cursor.fetchall()]

    conn.close()
    return sources


class ImportAnalyzeRequest(BaseModel):
    """导入分析请求"""
    content: str = Field(..., description="要分析的内容")
    auto_save: bool = Field(default=False, description="是否自动保存")


@router.post("/import/analyze")
async def import_analyze(request: ImportAnalyzeRequest):
    """
    智能分析导入内容，自动分类为四色卡片
    
    Args:
        request: 包含content和auto_save的请求体
        
    Returns:
        分析结果，包含分类后的卡片列表
    """
    try:
        content = request.content
        cards = []
        
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
        
        if not paragraphs:
            paragraphs = [p.strip() for p in content.split('\n') if p.strip()]
        
        for idx, para in enumerate(paragraphs):
            if len(para) < 10:
                continue
            
            # 过滤系统提示泄露内容（NPU模型有时会输出内部系统信息）
            leak_patterns = ['System Information', 'Template', 'DeepSeek', 'You are a helpful',
                           'instruction', 'prompt', 'special token', '#### Instruction',
                           '#### Response', 'system', '<|assistant', '<|end']
            if any(pat.lower() in para.lower() for pat in leak_patterns):
                continue
            # 过滤Markdown格式化的编号标题（如 "5. **xxx**:"）且内容过短
            if re.match(r'^\d+\.\s+\*\*', para) and len(para) < 80:
                continue
            
            lower_para = para.lower()
            
            card_type = 'blue'
            confidence = 0.5
            
            if any(kw in lower_para for kw in ['定义', '概念', '原理', '理论', '什么是']):
                card_type = 'blue'
                confidence = 0.8
            elif any(kw in lower_para for kw in ['关联', '相关', '连接', '对比', '区别']):
                card_type = 'green'
                confidence = 0.8
            elif any(kw in lower_para for kw in ['来源', '参考', '引用', 'http', 'www']):
                card_type = 'yellow'
                confidence = 0.8
            elif any(kw in lower_para for kw in ['关键词', '标签', '索引', '注意', '重要']):
                card_type = 'red'
                confidence = 0.8
            
            title = para[:50] + '...' if len(para) > 50 else para
            if '\n' in title:
                title = title.split('\n')[0]
            
            card = {
                'title': title,
                'content': para,
                'card_type': card_type,
                'confidence': confidence,
                'address': f"{card_type.upper()}{idx + 1}"
            }
            cards.append(card)
        
        if request.auto_save:
            conn = db_manager.get_connection()
            cursor = conn.cursor()
            for card in cards:
                cursor.execute('''
                    INSERT INTO knowledge_cards (card_type, title, content, category)
                    VALUES (?, ?, ?, ?)
                ''', (card['card_type'], card['title'], card['content'], card['card_type']))
            conn.commit()
            conn.close()
        
        return {
            'success': True,
            'cards': cards,
            'total': len(cards)
        }
        
    except Exception as e:
        logger.error(f"导入分析失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"分析失败: {str(e)}")


@router.post("/import/file")
async def import_file(file: UploadFile = File(...)):
    """
    导入文件并解析为知识卡片

    支持格式：PDF, TXT, MD, DOCX, XLSX, 图片
    """
    tmp_path = None
    source_file_id = None
    stored_file_path = None
    try:
        # 获取文件扩展名
        filename = file.filename or ""
        file_ext = os.path.splitext(filename)[1].lower()

        # 如果没有文件扩展名，尝试从content_type推断
        if not file_ext and file.content_type:
            content_type = file.content_type.lower()
            if 'pdf' in content_type:
                file_ext = '.pdf'
            elif 'text' in content_type or 'markdown' in content_type:
                file_ext = '.txt'
            elif 'word' in content_type or 'document' in content_type:
                file_ext = '.docx'
            elif 'excel' in content_type or 'spreadsheet' in content_type:
                file_ext = '.xlsx'
            elif 'image' in content_type:
                # Try to determine specific image format from content
                # Parse image subtype from content_type
                img_subtype = content_type.split('/')[-1].replace('jpeg', 'jpg')
                file_ext = f'.{img_subtype}' if img_subtype in ('jpg','png','gif','bmp','webp') else '.jpg'  # default to jpg

        # 如果仍然没有扩展名，尝试从文件内容检测（读取前几个字节）
        if not file_ext:
            # 保存上传的文件到临时目录以检测类型
            with tempfile.NamedTemporaryFile(delete=False) as tmp_check:
                content_preview = await file.read(1024)  # Read first 1024 bytes
                tmp_check.write(content_preview)
                tmp_check_path = tmp_check.name

            try:
                # Check file magic numbers
                with open(tmp_check_path, 'rb') as f:
                    header = f.read(8)

                if header.startswith(b'%PDF'):
                    file_ext = '.pdf'
                elif header.startswith((b'\xff\xd8\xff', b'\x89PNG', b'GIF8', b'BM')):
                    # Image formats
                    if header.startswith(b'\xff\xd8\xff'):
                        file_ext = '.jpg'
                    elif header.startswith(b'\x89PNG'):
                        file_ext = '.png'
                    elif header.startswith(b'GIF8'):
                        file_ext = '.gif'
                    elif header.startswith(b'BM'):
                        file_ext = '.bmp'
                elif b'\x00\x00\x00\x0c' in header[:4] or b'ftyp' in header:
                    # MP4/HEIC等格式
                    pass  # Unknown format, will be handled below
                else:
                    # Assume text file if no binary signature
                    file_ext = '.txt'  # fallback to text
            except Exception as e:
                logger.warning(f"无法检测文件类型: {e}")
                file_ext = '.txt'  # fallback to text
            finally:
                os.unlink(tmp_check_path)
                # Reset file pointer since we read some content
                await file.seek(0)

        # 验证支持的文件格式
        supported_extensions = {'.pdf', '.txt', '.md', '.docx', '.doc', '.xlsx', '.xls',
                              '.pptx', '.ppt',
                              '.mp3', '.wav', '.flac',
                              '.zip', '.rar', '.7z',
                              '.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp'}
        if file_ext not in supported_extensions:
            raise HTTPException(status_code=400, detail=f"请确保文件有正确的扩展名。检测到的格式: {file_ext}。支持的格式: {', '.join(sorted(supported_extensions))}")

        # 保存上传的文件到临时目录
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
            content = await file.read()
            if not content:
                raise HTTPException(status_code=400, detail="上传的文件为空")
            tmp.write(content)
            tmp_path = tmp.name

        # 生成源文件唯一标识符并保存到永久存储
        content_hash = hashlib.sha256(content).hexdigest()
        source_file_id = f"sf_{datetime.now().strftime('%Y%m%d')}_{content_hash[:12]}"
        
        # 创建源文件存储目录
        project_root = Path(__file__).parent.parent
        source_files_dir = project_root / "data" / "source_files"
        source_files_dir.mkdir(parents=True, exist_ok=True)
        
        # 保存文件到永久存储
        stored_file_path = source_files_dir / f"{source_file_id}{file_ext}"
        with open(stored_file_path, 'wb') as f:
            f.write(content)
        
        # 插入源文件记录到数据库
        conn = db_manager.get_connection()
        cursor = conn.cursor()
        try:
            # 检查是否已存在相同 hash 的记录
            cursor.execute("SELECT source_file_id FROM source_files WHERE content_hash = ?", (content_hash,))
            existing = cursor.fetchone()
            if existing:
                # 文件已存在，使用已有的 source_file_id
                source_file_id = existing[0]
                logger.info(f"文件已存在，使用已有记录: {source_file_id}")
            else:
                # 插入新记录
                cursor.execute('''
                    INSERT INTO source_files (source_file_id, original_name, stored_path, file_type, file_size, content_hash)
                    VALUES (?, ?, ?, ?, ?, ?)
                ''', (source_file_id, filename, str(stored_file_path), file_ext.lstrip('.'), len(content), content_hash))
                conn.commit()
        except Exception as e:
            logger.warning(f"插入源文件记录失败: {e}")
        conn.close()

        extracted_text = ""

        # 根据文件类型解析
        if file_ext == '.pdf':
            from tools.pdf_processor import SimplePDFProcessor
            processor = SimplePDFProcessor()
            result = processor.extract_text(tmp_path)
            extracted_text = result.get('full_text', '')

        elif file_ext in ['.txt', '.md']:
            with open(tmp_path, 'r', encoding='utf-8') as f:
                extracted_text = f.read()

        elif file_ext in ['.docx', '.doc']:
            try:
                from docx import Document
                doc = Document(tmp_path)
                extracted_text = '\n'.join([p.text for p in doc.paragraphs if p.text])
            except ImportError:
                raise HTTPException(status_code=503, detail="Word解析未安装，请运行: pip install python-docx")

        elif file_ext in ['.xlsx', '.xls']:
            try:
                import pandas as pd
                df = pd.read_excel(tmp_path)
                extracted_text = df.to_string()
            except ImportError:
                raise HTTPException(status_code=503, detail="Excel解析未安装，请运行: pip install pandas openpyxl")

        elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp']:
            # 图片需要视觉模型分析
            try:
                import httpx
                async with httpx.AsyncClient(timeout=60.0) as client:
                    # 优先使用项目自身的视觉分析API
                    with open(tmp_path, 'rb') as img_file:
                        response = await client.post(
                            "http://127.0.0.1:8000/api/vision/analyze",
                            files={"file": (filename, img_file, file.content_type or "image/jpeg")},
                            data={"question": "请提取图片中的所有文字和关键信息"}
                        )
                    if response.status_code == 200:
                        result = response.json()
                        if result.get("success"):
                            extracted_text = result.get("description", result.get("analysis", ""))
                            if result.get("facts"):
                                extracted_text += "\n\n" + "\n".join(result["facts"])
                        else:
                            extracted_text = result.get("analysis", "")
                    else:
                        logger.warning(f"视觉分析API返回 {response.status_code}")
                        extracted_text = ""
            except Exception as e:
                logger.warning(f"图片视觉分析失败: {e}")
                extracted_text = ""

        # ========== 新增格式支持 ==========
        elif file_ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(tmp_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                extracted_text += para.text + '\n'
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    extracted_text += cell.text + '\t'
                                extracted_text += '\n'
                    extracted_text += '\n---\n'
            except ImportError:
                raise HTTPException(status_code=503, detail="PPT解析未安装，请运行: pip install python-pptx")

        elif file_ext == '.ppt':
            # 旧版PPT，尝试通过LibreOffice转换
            try:
                import subprocess
                lo_paths = [
                    r'C:\Program Files\LibreOffice\program\soffice.exe',
                    r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
                ]
                soffice = None
                for p in lo_paths:
                    if os.path.exists(p):
                        soffice = p
                        break
                if not soffice:
                    soffice = 'soffice'
                pptx_path = tmp_path + 'x'  # .ppt -> .pptx
                result = subprocess.run(
                    [soffice, '--headless', '--convert-to', 'pptx', '--outdir', os.path.dirname(pptx_path), tmp_path],
                    capture_output=True, text=True, timeout=60
                )
                if os.path.exists(pptx_path):
                    from pptx import Presentation
                    prs = Presentation(pptx_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    extracted_text += para.text + '\n'
                        extracted_text += '\n---\n'
                    os.unlink(pptx_path)
                else:
                    extracted_text = f"[PPT文件转换失败]: {result.stderr[:200]}"
            except Exception as e:
                extracted_text = f"[PPT文件无法解析，建议转换为pptx格式]: {e}"

        elif file_ext in ['.mp3', '.wav', '.flac']:
            # 音频文件 -> STT语音识别转文字
            try:
                from services.speech_service import stt_transcribe_audio
                result = stt_transcribe_audio(str(tmp_path))
                if isinstance(result, dict):
                    extracted_text = result.get('text', '')
                else:
                    extracted_text = str(result)
            except ImportError:
                extracted_text = "[音频STT服务未安装，请确保已安装faster-whisper]"
            except Exception as e:
                logger.warning(f"音频STT识别失败: {e}")
                extracted_text = "[音频转文字失败]"

        elif file_ext in ['.zip', '.rar', '.7z']:
            # 压缩包：提取所有文本并合并
            import shutil
            extract_dir = tempfile.mkdtemp()
            try:
                all_texts = []
                if file_ext == '.zip':
                    import zipfile
                    with zipfile.ZipFile(tmp_path, 'r') as zf:
                        zf.extractall(extract_dir)
                elif file_ext == '.7z':
                    try:
                        import py7zr
                        with py7zr.SevenZipFile(tmp_path, 'r') as szf:
                            szf.extractall(extract_dir)
                    except ImportError:
                        raise HTTPException(status_code=503, detail="7z解析需要py7zr: pip install py7zr")
                elif file_ext == '.rar':
                    try:
                        import subprocess
                        # 尝试系统命令unrar
                        subprocess.run(['unrar', 'x', '-o+', tmp_path, extract_dir + os.sep],
                                      capture_output=True, timeout=60)
                    except Exception:
                        try:
                            import patoolib
                            patoolib.extract_archive(tmp_path, outdir=extract_dir)
                        except ImportError:
                            raise HTTPException(status_code=503, detail="RAR解析需要安装unrar或patool")

                # 递归读取提取出的文本文件
                for root, dirs, files in os.walk(extract_dir):
                    for fname in files:
                        fpath = os.path.join(root, fname)
                        fname_lower = fname.lower()
                        if fname_lower.endswith(('.txt', '.md', '.py', '.json', '.xml', '.html', '.csv')):
                            try:
                                with open(fpath, 'r', encoding='utf-8') as f:
                                    all_texts.append(f"### 文件: {fname}\n{f.read()}")
                            except Exception:
                                pass
                        elif fname_lower.endswith('.pdf'):
                            try:
                                from tools.pdf_processor import SimplePDFProcessor
                                proc = SimplePDFProcessor()
                                res = proc.extract_text(fpath)
                                all_texts.append(f"### 文件: {fname}\n{res.get('full_text', '')}")
                            except Exception:
                                pass
                        elif fname_lower.endswith(('.xlsx', '.xls')):
                            try:
                                import pandas as pd
                                df = pd.read_excel(fpath)
                                all_texts.append(f"### 文件: {fname}\n{df.to_string()}")
                            except Exception:
                                pass
                        elif fname_lower.endswith(('.docx', '.doc')):
                            try:
                                from docx import Document
                                doc = Document(fpath)
                                all_texts.append(f"### 文件: {fname}\n" + '\n'.join([p.text for p in doc.paragraphs if p.text]))
                            except Exception:
                                pass
                        elif fname_lower.endswith(('.pptx', '.ppt')):
                            try:
                                from pptx import Presentation
                                prs = Presentation(fpath)
                                ppt_text = ""
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if shape.has_text_frame:
                                            for p in shape.text_frame.paragraphs:
                                                ppt_text += p.text + '\n'
                                all_texts.append(f"### 文件: {fname}\n{ppt_text}")
                            except Exception:
                                pass

                extracted_text = '\n\n'.join(all_texts) if all_texts else f"[压缩包 {filename} 内未找到可解析的文本文件]"
            finally:
                shutil.rmtree(extract_dir, ignore_errors=True)

        # ========== 批量导入：直接用关键词规则，秒速完成 ==========
        # （8智能体调用LLM有超时等待，不适合批量导入）
        cards = []
        if extracted_text and len(extracted_text.strip()) > 20:
            paragraphs = [p.strip() for p in extracted_text.split('\n\n') if p.strip()]
            if not paragraphs:
                paragraphs = [p.strip() for p in extracted_text.split('\n') if p.strip()]

            for idx, para in enumerate(paragraphs):
                if len(para) < 10:
                    continue
                leak_patterns = ['System Information', 'Template', 'DeepSeek', 'You are a helpful',
                               'instruction', 'prompt', 'special token', '#### Instruction',
                               '#### Response', 'system', '<|assistant', '<|end']
                if any(pat.lower() in para.lower() for pat in leak_patterns):
                    continue
                if re.match(r'^\d+\.\s+\*\*', para) and len(para) < 80:
                    continue

                lower_para = para.lower()
                card_type = 'blue'
                confidence = 0.5
                if any(kw in lower_para for kw in ['定义', '概念', '原理', '理论', '什么是']):
                    card_type = 'blue'; confidence = 0.8
                elif any(kw in lower_para for kw in ['关联', '相关', '连接', '对比', '区别']):
                    card_type = 'green'; confidence = 0.8
                elif any(kw in lower_para for kw in ['来源', '参考', '引用', 'http', 'www']):
                    card_type = 'yellow'; confidence = 0.8
                elif any(kw in lower_para for kw in ['关键词', '标签', '索引', '注意', '重要']):
                    card_type = 'red'; confidence = 0.8

                title = para[:50] + '...' if len(para) > 50 else para
                if '\n' in title:
                    title = title.split('\n')[0]

                cards.append({
                    'title': title,
                    'content': para,
                    'card_type': card_type,
                    'confidence': confidence,
                    'address': f"{card_type.upper()}{idx + 1}"
                })

        saved_count = 0
        card_ids = []
        if cards:
            try:
                conn = db_manager.get_connection()
                cursor = conn.cursor()
                for idx, card in enumerate(cards):
                    category_map = {
                        'blue': '事实',
                        'green': '解释',
                        'yellow': '风险',
                        'red': '行动'
                    }
                    category = category_map.get(card['card_type'], '事实')
                    cursor.execute('''
                        INSERT INTO knowledge_cards (card_type, title, content, category, project_id)
                        VALUES (?, ?, ?, ?, ?)
                    ''', (
                        card['card_type'],
                        card['title'],
                        card['content'],
                        category,
                        None
                    ))
                    card_id = cursor.lastrowid
                    card_ids.append(card_id)
                    
                    # 插入卡片-源文件关联记录
                    if source_file_id:
                        location = f"第{idx + 1}段"
                        try:
                            cursor.execute('''
                                INSERT INTO card_source_files (source_file_id, card_id, location_in_source)
                                VALUES (?, ?, ?)
                            ''', (source_file_id, card_id, location))
                        except Exception as e:
                            logger.warning(f"插入卡片源文件关联失败: {e}")
                    
                    saved_count += 1
                conn.commit()
                conn.close()
                logger.info(f"成功保存 {saved_count} 张知识卡片")
            except Exception as e:
                logger.error(f"保存知识卡片失败: {e}")

        return {
            'success': True,
            'filename': filename,
            'file_type': file_ext,
            'extracted_length': len(extracted_text),
            'cards': cards,
            'total': len(cards),
            'saved': saved_count,
            'source_file_id': source_file_id  # 返回源文件ID供前端使用
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"文件导入失败: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"导入失败: {str(e)}")
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except Exception:
                pass


@router.post("/import/batch")
async def import_files_batch(files: List[UploadFile] = File(...)):
    """
    批量导入多个文件并解析为知识卡片

    并行处理多个文件的导入，返回每个文件的处理结果汇总
    """
    if not files:
        raise HTTPException(status_code=400, detail="请至少上传一个文件")

    results = []
    total_cards = 0
    success_count = 0
    fail_count = 0

    # 逐个处理（保持原导入逻辑，使用8智能体增强卡片生成）
    for idx, file in enumerate(files):
        try:
            # 复用单文件导入逻辑
            file_result = await import_file(file=file)
            results.append({
                "index": idx,
                "filename": file.filename or f"file_{idx}",
                "success": True,
                "file_type": file_result.get("file_type", ""),
                "cards_count": file_result.get("total", 0),
                "extracted_length": file_result.get("extracted_length", 0),
                "source_file_id": file_result.get("source_file_id", ""),
                "error": None
            })
            total_cards += file_result.get("total", 0)
            success_count += 1
        except HTTPException as e:
            fail_count += 1
            results.append({
                "index": idx,
                "filename": file.filename or f"file_{idx}",
                "success": False,
                "error": e.detail,
                "file_type": "",
                "cards_count": 0,
                "extracted_length": 0,
                "source_file_id": None
            })
        except Exception as e:
            fail_count += 1
            results.append({
                "index": idx,
                "filename": file.filename or f"file_{idx}",
                "success": False,
                "error": str(e),
                "file_type": "",
                "cards_count": 0,
                "extracted_length": 0,
                "source_file_id": None
            })

    return {
        "success": success_count > 0,
        "total": len(files),
        "success_count": success_count,
        "fail_count": fail_count,
        "total_cards": total_cards,
        "results": results
    }


@router.post("/import")
async def import_knowledge(html_dir: str):
    """
    导入知识库

    Args:
        html_dir: HTML 文件目录

    Returns:
        导入结果
    """
    try:
        # 添加项目路径到 sys.path
        import sys
        from pathlib import Path
        project_root = Path(__file__).parent.parent.parent
        if str(project_root) not in sys.path:
            sys.path.insert(0, str(project_root))

        # 动态导入批量导入工具
        import importlib
        import os

        # 设置 PYTHONPATH 环境变量
        os.environ['PYTHONPATH'] = str(project_root)

        # 导入批量导入模块
        spec = importlib.util.spec_from_file_location(
            "import_knowledge_batch",
            str(project_root / "backend" / "tools" / "import_knowledge_batch.py")
        )
        module = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(module)

        # 执行批量导入
        module.batch_import(html_dir)

        return {
            "success": True,
            "message": "知识库导入成功",
            "html_dir": html_dir
        }
    except Exception as e:
        import traceback
        logger.error(f"导入失败: {e}\n{traceback.format_exc()}")
        raise HTTPException(status_code=400, detail=f"导入失败: {str(e)}")


# ========== 专题关联接口 ==========

class CardTopicLink(BaseModel):
    """卡片专题关联"""
    card_id: int
    topic_id: int


@router.post("/cards/link-topic")
async def link_card_to_topic(link: CardTopicLink):
    """将卡片关联到专题 — 同时写入 project_id 和 topic_id 以保证兼容"""
    try:
        conn = db_manager.get_connection()
        cursor = conn.cursor()
        cursor.execute("""
            UPDATE knowledge_cards 
            SET project_id = ?, topic_id = ?, updated_at = CURRENT_TIMESTAMP
            WHERE id = ?
        """, (link.topic_id, link.topic_id, link.card_id))
        conn.commit()
        conn.close()
        return {"success": True, "message": f"卡片 {link.card_id} 已关联到专题 {link.topic_id}"}
    except Exception as e:
        logger.error(f"关联失败: {e}")
        raise HTTPException(status_code=400, detail=str(e))


@router.get("/cards/by-topic/{topic_id}")
async def get_cards_by_topic(topic_id: int):
    """获取专题下的所有卡片 — 同时匹配 project_id 和 topic_id"""
    try:
        conn = db_manager.get_connection()
        cursor = conn.cursor()
        cursor.execute("""
            SELECT id, title, content, card_type, category, project_id, topic_id, created_at
            FROM knowledge_cards 
            WHERE project_id = ? OR topic_id = ?
            ORDER BY created_at DESC
        """, (topic_id, topic_id))
        cards = [dict(row) for row in cursor.fetchall()]
        conn.close()
        return {"topic_id": topic_id, "cards": cards, "count": len(cards)}
    except Exception as e:
        logger.error(f"查询失败: {e}")
        raise HTTPException(status_code=400, detail=str(e))


@router.get("/cards/{card_id}/topics")
async def get_card_topics(card_id: int):
    """获取卡片关联的专题 — 优先返回 project_id，兼容 topic_id"""
    try:
        conn = db_manager.get_connection()
        cursor = conn.cursor()
        cursor.execute("""
            SELECT project_id, topic_id FROM knowledge_cards WHERE id = ?
        """, (card_id,))
        row = cursor.fetchone()
        conn.close()
        if row:
            project_id = row[0] or row[1]  # 优先 project_id，兼容 topic_id
            return {"card_id": card_id, "topic_id": project_id}
        return {"card_id": card_id, "topic_id": None}
    except Exception as e:
        logger.error(f"查询失败: {e}")
        raise HTTPException(status_code=400, detail=str(e))


# ==================== 源文件溯源功能 ====================

@router.get("/cards/{card_id}/source-file")
async def get_card_source_file(card_id: int):
    """
    获取卡片对应的源文件信息
    """
    try:
        conn = db_manager.get_connection()
        cursor = conn.cursor()
        
        # 查询卡片关联的源文件
        cursor.execute("""
            SELECT sf.id, sf.source_file_id, sf.original_name, sf.stored_path, 
                   sf.file_type, sf.file_size, sf.created_at, csf.location_in_source
            FROM card_source_files csf
            JOIN source_files sf ON csf.source_file_id = sf.source_file_id
            WHERE csf.card_id = ?
        """, (card_id,))
        
        row = cursor.fetchone()
        conn.close()
        
        if not row:
            return {"has_source": False, "message": "该卡片无溯源信息（非文件导入）"}
        
        return {
            "has_source": True,
            "source_file_id": row[1],
            "original_name": row[2],
            "stored_path": row[3],
            "file_type": row[4],
            "file_size": row[5],
            "created_at": row[6],
            "location_in_source": row[7]
        }
    except Exception as e:
        logger.error(f"获取源文件失败: {e}")
        raise HTTPException(status_code=400, detail=str(e))


@router.get("/source-files/{source_file_id}/cards")
async def get_source_file_cards(source_file_id: str):
    """
    获取某源文件生成的所有卡片
    """
    try:
        conn = db_manager.get_connection()
        cursor = conn.cursor()
        
        # 验证源文件是否存在
        cursor.execute("SELECT original_name, stored_path FROM source_files WHERE source_file_id = ?", (source_file_id,))
        sf_row = cursor.fetchone()
        if not sf_row:
            conn.close()
            raise HTTPException(status_code=404, detail="源文件不存在")
        
        # 获取关联的卡片
        cursor.execute("""
            SELECT kc.id, kc.title, kc.content, kc.card_type, kc.category, 
                   csf.location_in_source, kc.created_at
            FROM card_source_files csf
            JOIN knowledge_cards kc ON csf.card_id = kc.id
            WHERE csf.source_file_id = ?
            ORDER BY kc.created_at DESC
        """, (source_file_id,))
        
        cards = [dict(row) for row in cursor.fetchall()]
        conn.close()
        
        return {
            "source_file_id": source_file_id,
            "original_name": sf_row[0],
            "stored_path": sf_row[1],
            "cards": cards,
            "total": len(cards)
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"获取源文件卡片失败: {e}")
        raise HTTPException(status_code=400, detail=str(e))


@router.get("/source-files/{source_file_id}/download")
async def download_source_file(source_file_id: str):
    """
    下载源文件
    """
    try:
        conn = db_manager.get_connection()
        cursor = conn.cursor()
        
        cursor.execute("""
            SELECT original_name, stored_path FROM source_files WHERE source_file_id = ?
        """, (source_file_id,))
        
        row = cursor.fetchone()
        conn.close()
        
        if not row:
            raise HTTPException(status_code=404, detail="源文件不存在")
        
        original_name, stored_path = row
        
        if not os.path.exists(stored_path):
            raise HTTPException(status_code=404, detail="源文件已丢失")
        
        return FileResponse(
            path=stored_path,
            filename=original_name,
            media_type='application/octet-stream'
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"下载源文件失败: {e}")
        raise HTTPException(status_code=400, detail=str(e))


# 可视化图谱状态 API（持久化节点/连线/分类）
@router.get("/graph/state")
async def get_graph_state():
    """加载已保存的图谱状态"""
    try:
        state = db_manager.load_graph_state("default")
        if state:
            logger.info(f"[GraphState] 加载图谱: {len(state['nodes'])} 节点")
            return state
        return None
    except Exception as e:
        logger.error(f"[GraphState] 加载失败: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/graph/state")
async def save_graph_state(req: Request):
    """保存图谱状态（节点、连线、分类）"""
    try:
        body = await req.json()
        name = body.get("name", "default")
        nodes = body.get("nodes", [])
        links = body.get("links", [])
        categories = body.get("categories", [])
        success = db_manager.save_graph_state(name, nodes, links, categories)
        if success:
            return {"status": "ok", "nodes": len(nodes), "links": len(links)}
        raise HTTPException(status_code=500, detail="保存失败")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"[GraphState] 保存失败: {e}")
        raise HTTPException(status_code=400, detail=str(e))
