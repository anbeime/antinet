"""
Card → Direct PPTX Renderer using python-pptx (bypasses buggy SVG→DrawingML pipeline)
Dark-tech theme with gold accents, matching the Qualcomm Snapdragon competition style.
"""

import re, logging
from pathlib import Path
from typing import Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

SLIDE_W = 13.333
SLIDE_H = 7.5

# ── Color Palette (Dark Tech + Gold) ──
BG_DARK = RGBColor(0x0A, 0x0E, 0x1A)
BG_CARD = RGBColor(0x12, 0x1A, 0x2E)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
GOLD_BRIGHT = RGBColor(0xF5, 0xC8, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
MED_GRAY = RGBColor(0x7A, 0x84, 0x9A)
BLUE = RGBColor(0x4A, 0x9E, 0xF5)
GREEN = RGBColor(0x4E, 0xC9, 0x8B)
CARD_BLUE = RGBColor(0x3B, 0x82, 0xF6)
CARD_GREEN = RGBColor(0x22, 0xC5, 0x5E)
CARD_YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
CARD_RED = RGBColor(0xEF, 0x44, 0x44)

_INVALID_XML_RE = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F\uD800-\uDFFF\uFDD0-\uFDEF\uFFFE\uFFFF]")

def _sanitize(text: str) -> str:
    if not text:
        return ""
    return _INVALID_XML_RE.sub("", text)

# ── Theme System ──

_THEMES = {
    "dark-tech-gold": {
        "bg": BG_DARK,
        "bg_card": BG_CARD,
        "text": WHITE,
        "text_light": LIGHT_GRAY,
        "text_muted": MED_GRAY,
        "accent": GOLD,
        "accent_bright": GOLD_BRIGHT,
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
        "card_blue": CARD_BLUE,
        "card_green": CARD_GREEN,
        "card_yellow": CARD_YELLOW,
        "card_red": CARD_RED,
    },
    "professional": {
        "bg": RGBColor(0xF8, 0xFA, 0xFC),
        "bg_card": RGBColor(0xEE, 0xF1, 0xF8),
        "text": RGBColor(0x1F, 0x29, 0x37),
        "text_light": RGBColor(0x4B, 0x55, 0x6B),
        "text_muted": RGBColor(0x6B, 0x7A, 0x99),
        "accent": RGBColor(0x1E, 0x3A, 0x8A),
        "accent_bright": RGBColor(0x3B, 0x82, 0xF6),
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
        "card_blue": CARD_BLUE,
        "card_green": CARD_GREEN,
        "card_yellow": CARD_YELLOW,
        "card_red": CARD_RED,
    },
}

CARD_LABELS = {
    "blue": "核心概念",
    "green": "深度解读",
    "yellow": "风险警示",
    "red": "行动方案",
}


def _resolve_theme(name: str) -> dict:
    raw = _THEMES.get(name, _THEMES["dark-tech-gold"])
    t = dict(raw)
    t["card_map"] = {
        "blue": t["card_blue"],
        "green": t["card_green"],
        "yellow": t["card_yellow"],
        "red": t["card_red"],
    }
    return t


def _card_type(card_type: str, colors: dict) -> RGBColor:
    return colors["card_map"].get(card_type, colors["card_blue"])


# ── Design Helpers ──

def _add_shape(slide, left, top, w, h, fill=None, line=None, line_w=None, shape_type=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w or 1)
    else:
        s.line.fill.background()
    return s


def _add_rounded_rect(slide, left, top, w, h, fill, line=None, line_w=1):
    return _add_shape(slide, left, top, w, h, fill, line, line_w, MSO_SHAPE.ROUNDED_RECTANGLE)


def _add_accent_line(slide, left, top, w, color):
    _add_shape(slide, left, top, w, 0.04, fill=color)


def _add_type_tag(slide, left, top, w, h, label, bg_clr, text_clr=WHITE):
    """Small rounded pill in top-right of card: colored bg + light text."""
    _add_rounded_rect(slide, left, top, w, h, fill=bg_clr, line=bg_clr, line_w=0.5)
    _add_textbox(
        slide, left, top, w, h, label,
        "Microsoft YaHei", 11, text_clr,
        bold=True, alignment=PP_ALIGN.CENTER,
    )
    return left + w


def _add_textbox(slide, left, top, width, height, text, font_name, font_size, color,
                 bold=False, alignment=PP_ALIGN.LEFT, auto_shrink=True, anchor=MSO_ANCHOR.TOP):
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True
    # Inject normAutofit into a:bodyPr for auto-shrink behavior
    if auto_shrink:
        from pptx.oxml.ns import qn
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            for child in list(bodyPr):
                if child.tag in (qn('a:spAutoFit'), qn('a:normAutofit')):
                    bodyPr.remove(child)
            autofit_elem = bodyPr.makeelement(qn('a:normAutofit'), {})
            bodyPr.append(autofit_elem)
    p = tf.paragraphs[0]
    p.text = _sanitize(text)
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return txbox


def _solid_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Slide Builders ──

def _build_cover_slide(prs, topic, colors, card_count: int = 0, author: str = "Antinet 智能知识管家"):
    """封面：标题 + 元信息行（日期/作者/数量） + 横向分隔线"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, colors["bg"])

    now_str = datetime.now().strftime("%Y年%m月%d日")

    # 大标题（彩色）
    _add_textbox(
        slide, 0.7, 1.6, SLIDE_W - 1.4, 1.4,
        topic, colors["title_font"], 54, colors["accent_bright"],
        bold=True,
    )

    # 横向装饰线
    _add_shape(slide, 0.7, 3.2, SLIDE_W - 1.4, 0.05, fill=colors["accent"])

    # 元信息行：日期 | 作者 | 卡片数量
    meta = f"生成日期: {now_str}  |  作者: {author}"
    if card_count > 0:
        meta += f"  |  卡片数量: {card_count}"
    _add_textbox(
        slide, 0.7, 3.45, SLIDE_W - 1.4, 0.5,
        meta, colors["body_font"], 16, colors["text_light"],
    )

    # 底部脚注
    _add_textbox(
        slide, 0.7, SLIDE_H - 0.6, SLIDE_W - 1.4, 0.4,
        f"由 知易 · 锦衣卫多智能体 生成于 {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        colors["body_font"], 10, colors["text_muted"],
        alignment=PP_ALIGN.CENTER,
    )


def _build_section_slide(slide, card_type, title, content, colors, slide_num, total_slides,
                         source: str = ""):
    """知识卡片风格：一页一张卡，圆角彩色描边 + 编号彩色标题 + 类型徽章 + 内容 + 底部来源/时间戳"""
    clr = _card_type(card_type, colors)
    body_font = colors["body_font"]
    title_font = colors["title_font"]

    _solid_bg(slide, colors["bg"])

    # 圆角彩色描边卡片（占满整页，留 0.5" 边距）
    _add_rounded_rect(
        slide, 0.4, 0.3, SLIDE_W - 0.8, SLIDE_H - 0.6,
        fill=colors["bg"], line=clr, line_w=2.0,
    )

    # 顶部右上角：类型徽章（彩色实心 + 白字小药丸）
    type_label = CARD_LABELS.get(card_type, "笔记")
    tag_w, tag_h = 1.4, 0.42
    _add_type_tag(
        slide, SLIDE_W - 0.4 - tag_w, 0.55, tag_w, tag_h,
        type_label, clr, WHITE,
    )

    # 编号彩色标题（左上角，28pt 粗体）
    _add_textbox(
        slide, 0.7, 0.45, SLIDE_W - 2.5, 0.7,
        f"{slide_num}. {_sanitize(title[:60])}",
        title_font, 28, clr, bold=True,
    )

    # 内容正文（占据卡片主体，16pt 跟随主题文字色；启用自动缩放防止溢出）
    fs = 16 if len(content) < 400 else (14 if len(content) < 800 else 12)
    _add_textbox(
        slide, 0.7, 1.5, SLIDE_W - 1.4, SLIDE_H - 2.6,
        _sanitize(content), body_font, fs, colors["text"],
        auto_shrink=True,
    )

    # 底部：左侧来源路径，右侧创建时间（小号灰色）
    now_str = datetime.now().strftime("%Y/%m/%d %H:%M:%S")
    if source:
        left_footer = f"地址: {source}"
    else:
        left_footer = ""
    _add_textbox(
        slide, 0.7, SLIDE_H - 0.7, SLIDE_W - 1.4, 0.3,
        left_footer, body_font, 9, colors["text_muted"],
    )
    _add_textbox(
        slide, 0.7, SLIDE_H - 0.7, SLIDE_W - 1.4, 0.3,
        f"创建时间: {now_str}", body_font, 9, colors["text_muted"],
        alignment=PP_ALIGN.RIGHT,
    )


def _build_summary_slide(prs, cards, colors):
    """分析报告总结：4 张大色卡（按 type 计数）+ 居中大号数字 + 标签 + 底部脚注"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, colors["bg"])

    # 顶部装饰线
    _add_shape(slide, 0, 0, SLIDE_W, 0.06, fill=colors["accent"])
    _add_shape(slide, 0, 0, 0.10, SLIDE_H, fill=colors["accent"])

    # 标题
    _add_textbox(
        slide, 0.6, 0.55, SLIDE_W - 1.2, 0.7,
        "分析报告总结", colors["title_font"], 32,
        colors["accent_bright"], bold=True,
    )
    _add_textbox(
        slide, 0.6, 1.25, SLIDE_W - 1.2, 0.4,
        f"共 {len(cards)} 张分析卡片 · 按类型分布",
        colors["body_font"], 14, colors["text_light"],
    )

    # 4 张大色卡横向排列
    type_order = [
        ("blue", "核心概念"),
        ("green", "深度解读"),
        ("yellow", "风险警示"),
        ("red", "行动方案"),
    ]
    counts = {t: 0 for t, _ in type_order}
    for c in cards:
        t = c.get("type", "blue")
        if t in counts:
            counts[t] += 1

    n = len(type_order)
    margin = 0.6
    gap = 0.4
    total_w = SLIDE_W - 2 * margin
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 3.6
    top = 2.0

    for i, (ctype, label) in enumerate(type_order):
        clr = _card_type(ctype, colors)
        left = margin + i * (card_w + gap)
        _add_rounded_rect(
            slide, left, top, card_w, card_h,
            fill=colors["bg"], line=clr, line_w=2.5,
        )
        cx = left + card_w / 2
        _add_textbox(
            slide, left, top + 0.7, card_w, 1.6,
            str(counts[ctype]),
            colors["title_font"], 72, clr,
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide, left, top + 2.4, card_w, 0.5,
            f"{label}卡片",
            colors["body_font"], 16, colors["text"],
            bold=True, alignment=PP_ALIGN.CENTER,
        )

    # 底部脚注
    now_str = datetime.now().strftime("%Y-%m-%d %H:%M")
    _add_textbox(
        slide, 0.6, SLIDE_H - 0.55, SLIDE_W - 1.2, 0.3,
        f"共 {len(cards)} 张分析卡片  生成于 {now_str}",
        colors["body_font"], 10, colors["text_muted"],
        alignment=PP_ALIGN.CENTER,
    )


# ── Content Overflow Protection ──

_MAX_CONTENT_CHARS = 800  # max chars before auto-shrink kicks in
_MAX_CONTENT_CHARS_HARD = 2000  # force-split content beyond this

def _split_long_content(title: str, content: str, card_type: str, source: str) -> list[dict]:
    """Split overly long content across multiple logical slides."""
    if len(content) <= _MAX_CONTENT_CHARS_HARD:
        return [{"title": title, "content": content, "type": card_type, "source": source}]

    # Split by paragraphs (double newline)
    paragraphs = [p for p in content.split('\n\n') if p.strip()]
    slides = []
    current_chunk = []
    current_len = 0
    for para in paragraphs:
        if current_len + len(para) > _MAX_CONTENT_CHARS_HARD and current_chunk:
            slides.append("\n\n".join(current_chunk))
            current_chunk = [para]
            current_len = len(para)
        else:
            current_chunk.append(para)
            current_len += len(para) + 2
    if current_chunk:
        slides.append("\n\n".join(current_chunk))

    result = []
    for i, chunk in enumerate(slides):
        slide_title = f"{title}（续 {i+1}）" if i > 0 else title
        result.append({"title": slide_title, "content": chunk, "type": card_type, "source": source})
    return result


# ── Public API ──

def generate_pptx_direct(
    topic: str,
    cards: list[dict],
    theme_name: str = "dark-tech-gold",
    output_path: Optional[str] = None,
    source: str = "",
) -> str:
    """
    Generate PPTX directly using python-pptx with knowledge-card style.
    One slide per card, plus a cover and a 4-color stats summary.
    """
    colors = _resolve_theme(theme_name)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _build_cover_slide(prs, topic, colors, card_count=len(cards))

    total_cards = len(cards)
    slide_index = 1
    expanded_cards = []
    for card in cards:
        sub_cards = _split_long_content(
            card.get("title", ""),
            card.get("content", ""),
            card.get("type", "blue"),
            card.get("source", source),
        )
        expanded_cards.extend(sub_cards)

    total_slides = len(expanded_cards) + 2  # +cover +summary
    for card in expanded_cards:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _build_section_slide(
            slide,
            card.get("type", "blue"),
            card.get("title", ""),
            card.get("content", ""),
            colors,
            slide_index,
            total_slides,
            source=card.get("source", source),
        )
        slide_index += 1

    _build_summary_slide(prs, expanded_cards, colors)

    if output_path is None:
        output_dir = Path("C:/D/zhiyi/generated")
        output_dir.mkdir(parents=True, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = str(output_dir / f"{topic}_{ts}.pptx")

    prs.save(output_path)
    logger.info(f"PPTX generated: {output_path} ({len(prs.slides)} slides)")
    return output_path
