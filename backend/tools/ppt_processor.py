"""
PPT 处理器（增强版）
提供 PowerPoint 文档生成和处理功能，包括从文本自动生成 PPT
"""
import logging
import re
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx 未安装，PPT 功能不可用")
    # 提供占位符类以避免 NameError
    class RGBColor:
        def __init__(self, *args, **kwargs):
            pass
    class Presentation:
        pass
    class Inches:
        def __init__(self, *args, **kwargs):
            pass
    class Pt:
        def __init__(self, *args, **kwargs):
            pass
    class PP_ALIGN:
        CENTER = None

logger = logging.getLogger(__name__)


def clean_markdown_for_ppt(text: str) -> str:
    """
    清理Markdown标记，转换为PPT可用的纯文本
    """
    if not text:
        return ""
    
    # 替换链接: [文字](url) -> 文字
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    
    # 移除图片: ![文字](url) -> 
    text = re.sub(r'!\[([^\]]*)\]\([^)]+\)', '', text)
    
    # 移除代码块
    text = re.sub(r'```[^`]*```', '', text)
    text = re.sub(r'```[^`]*```', '', text)
    
    # 行内代码 `code` -> code
    text = re.sub(r'`([^`]+)`', r'\1', text)
    
    # 加粗 **text** -> text
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    
    # 斜化 *text* -> text  
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    
    # 删除线 ~~text~~ -> text
    text = re.sub(r'~~([^~]+)~~', r'\1', text)
    
    # 清理多余空白
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = text.strip()
    
    return text


def parse_markdown_content(content: str) -> List[Dict[str, Any]]:
    """
    解析 Markdown 内容为幻灯片结构
    
    Args:
        content: Markdown 格式的文本内容
        
    Returns:
        幻灯片数据列表
    """
    slides = []
    current_slide = None
    
    lines = content.split('\n')
    
    for line in lines:
        line_stripped = line.strip()
        
        # 一级标题 - 新幻灯片标题
        if line_stripped.startswith('# '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'title': line_stripped[2:].strip(),
                'content': [],
                'type': 'title'
            }
        
        # 二级标题 - 新幻灯片或章节
        elif line_stripped.startswith('## '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'title': line_stripped[3:].strip(),
                'content': [],
                'type': 'content'
            }
        
        # 三级标题 - 内容标题
        elif line_stripped.startswith('### '):
            if current_slide is None:
                current_slide = {
                    'title': line_stripped[4:].strip(),
                    'content': [],
                    'type': 'content'
                }
            else:
                current_slide['content'].append({
                    'type': 'heading',
                    'text': line_stripped[4:].strip()
                })
        
        # 列表项
        elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
            if current_slide is None:
                current_slide = {
                    'title': '内容',
                    'content': [],
                    'type': 'content'
                }
            current_slide['content'].append({
                'type': 'bullet',
                'text': line_stripped[2:].strip()
            })
        
        # 编号列表
        elif re.match(r'^\d+\.\s', line_stripped):
            if current_slide is None:
                current_slide = {
                    'title': '内容',
                    'content': [],
                    'type': 'content'
                }
            text = re.sub(r'^\d+\.\s', '', line_stripped)
            current_slide['content'].append({
                'type': 'numbered',
                'text': text.strip()
            })
        
        # 普通段落
        elif line_stripped and not line_stripped.startswith('#'):
            if current_slide is None:
                current_slide = {
                    'title': '内容',
                    'content': [],
                    'type': 'content'
                }
            current_slide['content'].append({
                'type': 'paragraph',
                'text': line_stripped
            })
    
    # 添加最后一个幻灯片
    if current_slide:
        slides.append(current_slide)
    
    return slides


class PPTProcessor:
    """PPT 处理器类"""
    
    # 四色卡片颜色映射
    CARD_COLORS = {
        "fact": RGBColor(52, 152, 219),      # 蓝色 - 事实
        "interpret": RGBColor(46, 204, 113),  # 绿色 - 解释
        "risk": RGBColor(241, 196, 15),       # 黄色 - 风险
        "action": RGBColor(231, 76, 60)       # 红色 - 行动
    }
    
    CARD_NAMES = {
        "fact": "事实卡片",
        "interpret": "解释卡片",
        "risk": "风险卡片",
        "action": "行动卡片"
    }
    
    # 主题配色方案
    THEMES = {
        "professional": {
            "primary": RGBColor(28, 40, 51),      # 深蓝灰
            "secondary": RGBColor(52, 152, 219),  # 蓝色
            "accent": RGBColor(241, 196, 15),     # 金色
            "text": RGBColor(44, 62, 80),         # 深灰
            "background": RGBColor(236, 240, 241) # 浅灰
        },
        "creative": {
            "primary": RGBColor(155, 89, 182),    # 紫色
            "secondary": RGBColor(52, 152, 219),  # 蓝色
            "accent": RGBColor(230, 126, 34),     # 橙色
            "text": RGBColor(44, 62, 80),         # 深灰
            "background": RGBColor(236, 240, 241) # 浅灰
        },
        "minimal": {
            "primary": RGBColor(44, 62, 80),      # 深灰
            "secondary": RGBColor(149, 165, 166), # 中灰
            "accent": RGBColor(52, 152, 219),     # 蓝色
            "text": RGBColor(44, 62, 80),         # 深灰
            "background": RGBColor(255, 255, 255) # 白色
        }
    }
    
    def __init__(self):
        """初始化 PPT 处理器"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")
    
    def create_presentation(self, title: str = "Antinet 智能分析报告") -> Presentation:
        """
        创建新的演示文稿
        
        Args:
            title: 演示文稿标题
            
        Returns:
            Presentation 对象
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # 添加标题页
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        return prs
    
    def create_from_text(
        self,
        content: str,
        output_path: str,
        title: str = "演示文稿",
        theme: str = "professional"
    ) -> str:
        """
        从文本内容创建 PPT
        
        Args:
            content: 文本内容（支持 Markdown）
            output_path: 输出文件路径
            title: 演示文稿标题
            theme: 主题风格
            
        Returns:
            输出文件路径
        """
        try:
            slides_data = parse_markdown_content(content)
            
            if not slides_data:
                raise ValueError("无法从内容中解析出幻灯片")
            
            return self.create_presentation_from_slides(
                slides_data=slides_data,
                title=title,
                output_path=output_path,
                theme=theme
            )
            
        except Exception as e:
            logger.error(f"从文本生成 PPT 失败: {e}", exc_info=True)
            raise
    
    def create_presentation_from_slides(
        self,
        slides_data: List[Dict[str, Any]],
        title: str,
        output_path: str,
        theme: str = "professional"
    ) -> str:
        """
        从幻灯片数据创建 PPT
        
        Args:
            slides_data: 幻灯片数据列表
            title: 演示文稿标题
            output_path: 输出文件路径
            theme: 主题风格
            
        Returns:
            输出文件路径
        """
        try:
            prs = self.create_presentation(title)
            
            theme_colors = self.THEMES.get(theme, self.THEMES["professional"])
            
            for slide_data in slides_data:
                self._add_text_slide(prs, slide_data, theme_colors)
            
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
            
            logger.info(f"从幻灯片数据生成 PPT 成功: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"从幻灯片数据生成 PPT 失败: {e}", exc_info=True)
            raise
    
    def create_presentation_from_cards(
        self,
        cards: List[Dict[str, Any]],
        title: str,
        output_path: str,
        include_summary: bool = True
    ) -> str:
        """
        从卡片数据创建 PPT
        
        Args:
            cards: 卡片列表
            title: 演示文稿标题
            output_path: 输出文件路径
            include_summary: 是否包含总结页
            
        Returns:
            输出文件路径
        """
        return self.export_cards_to_ppt(
            cards=cards,
            output_path=output_path,
            title=title,
            include_summary=include_summary
        )
    
    def _add_text_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        theme_colors: Dict[str, RGBColor]
    ) -> None:
        """
        添加文本幻灯片
        
        Args:
            prs: Presentation 对象
            slide_data: 幻灯片数据
            theme_colors: 主题配色
        """
        if slide_data['type'] == 'title':
            # 标题页
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            
            if slide_data['content']:
                subtitle_text = '\n'.join([
                    item['text'] for item in slide_data['content']
                    if item['type'] == 'paragraph'
                ])
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = subtitle_text
        
        else:
            # 内容页 - 使用空白布局自定义
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # 添加标题
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = theme_colors['primary']
            
            # 添加内容
            if slide_data['content']:
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5),
                    Inches(9), Inches(5.5)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                for i, item in enumerate(slide_data['content']):
                    if i > 0:
                        p = content_frame.add_paragraph()
                    else:
                        p = content_frame.paragraphs[0]
                    
                    if item['type'] == 'heading':
                        p.text = item['text']
                        p.font.size = Pt(20)
                        p.font.bold = True
                        p.font.color.rgb = theme_colors['secondary']
                        p.space_before = Pt(12)
                    
                    elif item['type'] == 'bullet':
                        p.text = item['text']
                        p.level = 0
                        p.font.size = Pt(16)
                        p.font.color.rgb = theme_colors['text']
                    
                    elif item['type'] == 'numbered':
                        p.text = item['text']
                        p.level = 0
                        p.font.size = Pt(16)
                        p.font.color.rgb = theme_colors['text']
                    
                    elif item['type'] == 'paragraph':
                        p.text = item['text']
                        p.font.size = Pt(14)
                        p.font.color.rgb = theme_colors['text']
                        p.space_after = Pt(6)
    
    def add_card_slide(self, prs: Presentation, card: Dict[str, Any]) -> None:
        """
        添加卡片幻灯片
        
        Args:
            prs: Presentation 对象
            card: 卡片数据字典
        """
        # 使用空白布局
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        card_type = card.get("type", "fact")
        card_color = self.CARD_COLORS.get(card_type, RGBColor(128, 128, 128))
        card_name = self.CARD_NAMES.get(card_type, "卡片")
        
        # 添加卡片类型标题（左上角色块）
        type_box = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0.5), Inches(0.5),
            Inches(2), Inches(0.6)
        )
        type_box.fill.solid()
        type_box.fill.fore_color.rgb = card_color
        type_box.line.color.rgb = card_color
        
        type_text = type_box.text_frame
        type_text.text = card_name
        type_text.paragraphs[0].font.size = Pt(20)
        type_text.paragraphs[0].font.bold = True
        type_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        type_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 添加卡片标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = card.get("title", "无标题")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        
        # 添加卡片内容
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.3),
            Inches(9), Inches(4.2)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        content = card.get("content", "")
        if isinstance(content, list):
            content = "\n".join(f"• {item}" for item in content)
        
        content_frame.text = content
        content_frame.paragraphs[0].font.size = Pt(16)
        content_frame.paragraphs[0].font.color.rgb = RGBColor(52, 73, 94)
        content_frame.paragraphs[0].line_spacing = 1.5
        
        # 添加底部元数据
        if card.get("tags") or card.get("created_at"):
            meta_text = []
            if card.get("tags"):
                tags = card["tags"] if isinstance(card["tags"], list) else [card["tags"]]
                meta_text.append(f"标签: {', '.join(tags)}")
            if card.get("created_at"):
                meta_text.append(f"创建时间: {card['created_at']}")
            
            meta_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.8),
                Inches(9), Inches(0.4)
            )
            meta_frame = meta_box.text_frame
            meta_frame.text = " | ".join(meta_text)
            meta_frame.paragraphs[0].font.size = Pt(10)
            meta_frame.paragraphs[0].font.color.rgb = RGBColor(149, 165, 166)
    
    def add_summary_slide(self, prs: Presentation, summary: Dict[str, Any]) -> None:
        """
        添加总结幻灯片
        
        Args:
            prs: Presentation 对象
            summary: 总结数据字典
        """
        title_content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_content_layout)
        
        title = slide.shapes.title
        title.text = summary.get("title", "分析总结")
        
        content_box = slide.placeholders[1]
        text_frame = content_box.text_frame
        text_frame.clear()
        
        # 添加总结内容
        for key, value in summary.items():
            if key == "title":
                continue
            
            p = text_frame.add_paragraph()
            p.text = f"{key}: {value}"
            p.level = 0
            p.font.size = Pt(18)
    
    def add_chart_slide(self, prs: Presentation, title: str, chart_data: Dict[str, Any]) -> None:
        """
        添加图表幻灯片
        
        Args:
            prs: Presentation 对象
            title: 幻灯片标题
            chart_data: 图表数据
        """
        # 使用标题+内容布局
        title_content_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_content_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # 添加图表说明文本
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(4)
        )
        text_frame = content_box.text_frame
        text_frame.text = "图表数据:\n" + str(chart_data)
        text_frame.paragraphs[0].font.size = Pt(14)
    
    def export_cards_to_ppt(
        self,
        cards: List[Dict[str, Any]],
        output_path: str,
        title: str = "Antinet 四色卡片分析报告",
        include_summary: bool = True
    ) -> str:
        """
        将四色卡片导出为 PPT
        
        Args:
            cards: 卡片列表
            output_path: 输出文件路径
            title: 演示文稿标题
            include_summary: 是否包含总结页
            
        Returns:
            输出文件路径
        """
        try:
            # 创建演示文稿
            prs = self.create_presentation(title)
            
            # 按类型分组卡片
            cards_by_type = {
                "fact": [],
                "interpret": [],
                "risk": [],
                "action": []
            }
            
            for card in cards:
                card_type = card.get("type", "fact")
                if card_type in cards_by_type:
                    cards_by_type[card_type].append(card)
            
            # 添加卡片幻灯片（按类型顺序）
            for card_type in ["fact", "interpret", "risk", "action"]:
                type_cards = cards_by_type[card_type]
                if type_cards:
                    for card in type_cards:
                        self.add_card_slide(prs, card)
            
            # 添加总结页
            if include_summary:
                summary = {
                    "title": "分析总结",
                    "总卡片数": len(cards),
                    "事实卡片": len(cards_by_type["fact"]),
                    "解释卡片": len(cards_by_type["interpret"]),
                    "风险卡片": len(cards_by_type["risk"]),
                    "行动卡片": len(cards_by_type["action"])
                }
                self.add_summary_slide(prs, summary)
            
            # 保存文件
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
            
            logger.info(f"PPT 导出成功: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"PPT 导出失败: {e}", exc_info=True)
            raise
    
    def create_analysis_report(
        self,
        analysis_data: Dict[str, Any],
        output_path: str
    ) -> str:
        """
        创建完整的分析报告 PPT
        
        Args:
            analysis_data: 分析数据，包含 cards, charts, summary 等
            output_path: 输出文件路径
            
        Returns:
            输出文件路径
        """
        try:
            title = analysis_data.get("title", "Antinet 智能分析报告")
            prs = self.create_presentation(title)
            
            # 添加卡片
            cards = analysis_data.get("cards", [])
            for card in cards:
                self.add_card_slide(prs, card)
            
            # 添加图表
            charts = analysis_data.get("charts", [])
            for chart in charts:
                self.add_chart_slide(
                    prs,
                    chart.get("title", "数据图表"),
                    chart.get("data", {})
                )
            
            # 添加总结
            summary = analysis_data.get("summary")
            if summary:
                self.add_summary_slide(prs, summary)
            
            # 保存
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
            
            logger.info(f"分析报告 PPT 创建成功: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"创建分析报告失败: {e}", exc_info=True)
            raise


# 便捷函数
def export_cards_to_ppt(
    cards: List[Dict[str, Any]],
    output_path: str,
    title: str = "Antinet 四色卡片分析报告"
) -> str:
    """
    便捷函数：将卡片导出为 PPT
    
    Args:
        cards: 卡片列表
        output_path: 输出路径
        title: 标题
        
    Returns:
        输出文件路径
    """
    processor = PPTProcessor()
    return processor.export_cards_to_ppt(cards, output_path, title)


def create_analysis_ppt(
    analysis_data: Dict[str, Any],
    output_path: str
) -> str:
    """
    便捷函数：创建分析报告 PPT
    
    Args:
        analysis_data: 分析数据
        output_path: 输出路径
        
    Returns:
        输出文件路径
    """
    processor = PPTProcessor()
    return processor.create_analysis_report(analysis_data, output_path)
