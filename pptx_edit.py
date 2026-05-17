import sys
sys.path.insert(0, r'C:\D\zhiyi\venv_x64\Lib\site-packages')
from pptx import Presentation
import os

# Try Windows path
path = r'C:\D\决赛答辩PPT模板_正式版.pptx'
print('Exists:', os.path.exists(path))

prs = Presentation(path)
print(f'Slides: {len(prs.slides)}')
for i, layout in enumerate(prs.slide_layouts):
    print(f'Layout {i}: {layout.name}')
for i, slide in enumerate(prs.slides):
    print(f'\nSlide {i+1}: {slide.slide_layout.name}')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(f'  {shape.text[:100]}')