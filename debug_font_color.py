import requests, json, glob, os, io
from pptx import Presentation
from pptx.util import Emu

# Use latest generated PPTX
files = sorted(glob.glob(r'C:\D\zhiyi\generated\*.pptx'))
if files:
    prs = Presentation(files[-1])
    slide = prs.slides[0]
    for shape in slide.shapes:
        print(f"Shape: {shape.name}, type={shape.shape_type}, has_text={shape.has_text_frame}")
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    try:
                        clr_info = f"color.rgb={run.font.color.rgb}" if (run.font.color and run.font.color.rgb) else "no rgb"
                    except Exception as e:
                        clr_info = f"EXCEPTION: {e}"
                    print(f"  p{pi}r{ri}: text={run.text[:40]}, size={run.font.size}, bold={run.font.bold}, {clr_info}")
        print()
else:
    print("No PPTX files")
